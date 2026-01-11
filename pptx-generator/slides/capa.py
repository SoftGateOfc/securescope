import sys
import base64
from io import BytesIO
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from utils.processar_dados import fix_encoding


def add_base64_image(slide, base64_str, left, top, width, height):
    try:
        if ',' in base64_str:
            base64_str = base64_str.split(',')[1]

        img_data = base64.b64decode(base64_str)
        stream = BytesIO(img_data)
        slide.shapes.add_picture(
            stream, Inches(left), Inches(top),
            width=Inches(width), height=Inches(height)
        )
        return True
    except Exception as e:
        print(f"Erro imagem: {e}", file=sys.stderr)
        return False


def gerar_capa(pres, dados):
    print("Criando capa corporativa...", file=sys.stderr)
    slide = pres.slides.add_slide(pres.slide_layouts[6])

    # === IMAGEM DE FUNDO (primeiro elemento, fica atrás de tudo) === #
    try:
        with open('images/CAPA.png', 'rb') as f:
            img_data = f.read()
            stream = BytesIO(img_data)
            # Adiciona a imagem ocupando todo o slide
            slide.shapes.add_picture(
                stream, 
                Inches(0), 
                Inches(0),
                width=pres.slide_width,
                height=pres.slide_height
            )
    except Exception as e:
        print(f"Erro ao carregar imagem de fundo: {e}", file=sys.stderr)

    # === LOGO DA EMPRESA === #
    logo = dados.get("imagens", {}).get("logo_empresa")
    if logo:
        add_base64_image(slide, logo, 0.35, 0.35, 1.5, 0.92)

    # === TÍTULO PRINCIPAL === #
    title_left = Inches(4.5)
    title_top = Inches(1.5)
    title_w = Inches(6.5)
    title_h = Inches(1.8)

    title_box = slide.shapes.add_textbox(title_left, title_top, title_w, title_h)
    tf = title_box.text_frame
    tf.clear()

    # Primeira linha: ANÁLISE DE
    p1 = tf.paragraphs[0]
    p1.text = "ANÁLISE DE"
    p1.font.name = "Calibri"
    p1.font.size = Pt(50)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(255, 255, 255)

    # Segunda linha: RISCOS
    p2 = tf.add_paragraph()
    p2.text = "RISCOS"
    p2.font.name = "Calibri"
    p2.font.size = Pt(50)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(255, 255, 255)

    # === LINHA DECORATIVA DUPLA === #
    # Linha superior (fina)
    linha1 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        title_left,
        title_top + Inches(1.8),
        Inches(4.0),
        Inches(0.02)
    )
    linha1.fill.solid()
    linha1.fill.fore_color.rgb = RGBColor(255, 255, 255)
    linha1.line.fill.background()

    # Linha inferior (grossa)
    linha2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        title_left,
        title_top + Inches(1.88),
        Inches(4.0),
        Inches(0.05)
    )
    linha2.fill.solid()
    linha2.fill.fore_color.rgb = RGBColor(255, 255, 255)
    linha2.line.fill.background()

    # === INFORMAÇÕES DA EMPRESA === # 
    
    empresa = fix_encoding(dados.get("dados", {}).get("nome_empresa", ""))
    local = fix_encoding(dados.get("dados", {}).get("localizacao_analise", ""))
    
    info_top = title_top + Inches(2.05)
    
    # Nome da empresa
    if empresa:
        empresa_box = slide.shapes.add_textbox(
            title_left,
            info_top,
            Inches(4.0),
            Inches(0.5)
        )
        tf_emp = empresa_box.text_frame
        tf_emp.clear()
        tf_emp.word_wrap = True
        p_emp = tf_emp.paragraphs[0]
        p_emp.text = empresa.upper()
        p_emp.font.name = "Calibri"
        p_emp.font.size = Pt(22)
        p_emp.font.bold = True
        p_emp.font.color.rgb = RGBColor(255, 255, 255)

    # Localização
    if local:
        loc_box = slide.shapes.add_textbox(
            title_left,
            info_top + Inches(0.45),
            Inches(4.0),
            Inches(0.6)
        )
        tf_loc = loc_box.text_frame
        tf_loc.clear()
        tf_loc.word_wrap = True
        p_loc = tf_loc.paragraphs[0]
        p_loc.text = local.upper()
        p_loc.font.name = "Calibri Light"
        p_loc.font.size = Pt(16)
        p_loc.font.color.rgb = RGBColor(255, 255, 255)

    # === DATA (MÊS / ANO em duas linhas) === #
    agora = datetime.now()
    
    meses_pt = {
        'January': 'JANEIRO',
        'February': 'FEVEREIRO',
        'March': 'MARÇO',
        'April': 'ABRIL',
        'May': 'MAIO',
        'June': 'JUNHO',
        'July': 'JULHO',
        'August': 'AGOSTO',
        'September': 'SETEMBRO',
        'October': 'OUTUBRO',
        'November': 'NOVEMBRO',
        'December': 'DEZEMBRO'
    }
    
    mes_ingles = agora.strftime("%B")
    mes = meses_pt.get(mes_ingles, mes_ingles.upper())
    ano = agora.strftime("%Y")
    ano_espacado = " ".join(ano)
    
    data_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.7),
        Inches(2.0), Inches(0.6)
    )
    tf_d = data_box.text_frame
    tf_d.clear()
    
    # Primeira linha: MÊS (em português)
    p_mes = tf_d.paragraphs[0]
    p_mes.text = mes
    p_mes.font.name = "Calibri"
    p_mes.font.size = Pt(16)
    p_mes.font.bold = True
    p_mes.font.color.rgb = RGBColor(255, 255, 255)
    
    # Segunda linha: ANO (com espaçamento entre dígitos)
    p_ano = tf_d.add_paragraph()
    p_ano.text = ano_espacado
    p_ano.font.name = "Arial"
    p_ano.font.size = Pt(24)
    p_ano.font.bold = True
    p_ano.font.color.rgb = RGBColor(255, 255, 255)

    print("Capa corporativa criada!", file=sys.stderr)
    return slide