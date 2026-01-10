#!/usr/bin/env python3
# -*- coding: utf-8 -*-
import sys
import os
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


def add_local_image(slide, image_path, left, top, width, height):
    """Adiciona imagem local ao slide"""
    try:
        if os.path.exists(image_path):
            slide.shapes.add_picture(
                image_path,
                Inches(left),
                Inches(top),
                width=Inches(width),
                height=Inches(height)
            )
            return True
        else:
            print(f"⚠️ Imagem não encontrada: {image_path}", file=sys.stderr)
            return False
    except Exception as e:
        print(f"❌ Erro ao adicionar imagem: {e}", file=sys.stderr)
        return False

def tracar_linha_horizontal(slide, y, x_start, comprimento, cor_rgb=(70, 70, 70), espessura=0.5, circulo_final=True, texto_circulo= "0%",  titulo="Título",
    cor_titulo=(0, 0, 0),
    descricao="Descrição do indicador"):
    """
    Traça uma linha horizontal no slide
    
    Parâmetros:
    - slide: O slide onde a linha será adicionada
    - y: Altura da linha em inches (distância do topo)
    - x_start: Posição inicial da linha em inches (distância da esquerda)
    - comprimento: Comprimento da linha em inches
    - cor_rgb: Tupla RGB da cor (padrão: cinza 70,70,70)
    - espessura: Espessura da linha em pontos (padrão: 0.5)
    
    Exemplo de uso:
    tracar_linha_horizontal(slide, y=2.5, x_start=3.0, comprimento=4.5)
    tracar_linha_horizontal(slide, y=3.0, x_start=2.0, comprimento=6.0, cor_rgb=(255, 0, 0), espessura=1.0)
    """
    x_end = x_start + comprimento
    
    linha = slide.shapes.add_connector(
        1,  # Tipo 1 = linha reta
        Inches(x_start),
        Inches(y),
        Inches(x_end),
        Inches(y)
    )
    
    linha.line.color.rgb = RGBColor(*cor_rgb)
    linha.line.width = Pt(espessura)

    if circulo_final:
        circulo_tamanho = 0.08  # Tamanho do círculo em inches (bem pequeno)
        
        circulo = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x_end - circulo_tamanho/2),  # Centralizado no fim da linha
            Inches(y - circulo_tamanho/2),      # Centralizado verticalmente
            Inches(circulo_tamanho),
            Inches(circulo_tamanho)
        )
        
        # Preencher o círculo com a mesma cor da linha
        circulo.fill.solid()
        circulo.fill.fore_color.rgb = RGBColor(*cor_rgb)
        
        # Remover borda do círculo
        circulo.line.fill.background()

        diametro_grande = 1
        deslocamento_centro = 0.5  # distância do centro do círculo pequeno

        centro_pequeno_x = x_end
        centro_pequeno_y = y

        centro_grande_x = centro_pequeno_x + deslocamento_centro +0.1

        circulo_grande = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(centro_grande_x - diametro_grande / 2),
            Inches(centro_pequeno_y - diametro_grande / 2),
            Inches(diametro_grande),
            Inches(diametro_grande)
        )

        fill = circulo_grande.fill
        fill.gradient()
        fill.gradient_angle = 90  # Vertical

        # Stop 0: azul escuro (embaixo)
        stop_dark = fill.gradient_stops[0]
        stop_dark.position = 0.0
        stop_dark.color.rgb = RGBColor(10, 50, 100)

        # Stop 1: azul claro (em cima)
        stop_light = fill.gradient_stops[1]
        stop_light.position = 1.0
        stop_light.color.rgb = RGBColor(30, 115, 190)

        circulo_grande.line.fill.background()

         # TEXTO CENTRAL 
        tf = circulo_grande.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        tf.margin_left = Inches(0.02)
        tf.margin_right = Inches(0.02)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)

        p = tf.paragraphs[0]
        p.clear()
        p.alignment = PP_ALIGN.CENTER

        # Extrair número e símbolo %
        valor = str(texto_circulo).replace('%', '')

        # 🔢 RUN DO NÚMERO (18pt)
        run_numero = p.add_run()
        run_numero.text = valor
        run_numero.font.name = "Arial"
        run_numero.font.size = Pt(18)
        run_numero.font.bold = True
        run_numero.font.color.rgb = RGBColor(255, 255, 255)

        # 🔣 RUN DO %
        run_percent = p.add_run()
        run_percent.text = "%"
        run_percent.font.name = "Arial"
        run_percent.font.size = Pt(14)
        run_percent.font.bold = True
        run_percent.font.color.rgb = RGBColor(255, 255, 255)

         # 📝 BLOCO DE TEXTO (TÍTULO + DESCRIÇÃO)
        largura_texto = 2.5
        offset_texto = diametro_grande / 2 + 0.05

        caixa_texto = slide.shapes.add_textbox(
            Inches(centro_grande_x + offset_texto),
            Inches(y - 0.25),
            Inches(largura_texto),
            Inches(0.8)
        )

        tf_texto = caixa_texto.text_frame
        tf_texto.clear()

        # 🔹 TÍTULO
        p_titulo = tf_texto.paragraphs[0]
        p_titulo.text = titulo
        p_titulo.font.name = "Arial"
        p_titulo.font.size = Pt(14)
        p_titulo.font.bold = True
        p_titulo.font.color.rgb = RGBColor(*cor_titulo)

        # 🔸 DESCRIÇÃO
        p_desc = tf_texto.add_paragraph()
        p_desc.text = descricao
        p_desc.font.name = "Arial"
        p_desc.font.size = Pt(9)
        p_desc.font.color.rgb = RGBColor(70, 70, 70)

    
    return linha

def gerar_resumo_exec(pres, dados):
    percentuais = dados.get("dados_modelo", {}).get("porcentagem_pilar", {})
    """Gera o slide de Resumo Executivo (Página 5)"""
    print("Criando slide de Resumo Executivo...", file=sys.stderr)
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    tracar_linha_horizontal(slide, y=0.9, x_start=3.55, comprimento=3.6, texto_circulo= f"{percentuais.get('Tecnologia', 0)}%",titulo="TECNOLOGIA",
    cor_titulo=(10, 50, 100),
    descricao="Tecnologias atuais ampliam \nvigilância e pronta reação.")
    tracar_linha_horizontal(slide, y=1.8, x_start=4.4, comprimento=2.1,texto_circulo= f"{percentuais.get('Processos', 0)}%", titulo="PROCESSOS",
    cor_titulo=(10, 50, 100) ,
    descricao="Revisão contínua alinha fluxos \na padrões de controles.")
    tracar_linha_horizontal(slide, y=2.8, x_start=3.4, comprimento=2.6, texto_circulo= f"{percentuais.get('Pessoas', 0)}%", titulo="PESSOAS",
    cor_titulo=(10, 50, 100) ,
    descricao="Cultura de segurança internalizada \npor todos os times.")
    tracar_linha_horizontal(slide, y=3.85, x_start=4.4, comprimento=2.1,texto_circulo= f"{percentuais.get('Gestao', 0)}%", titulo="GESTÃO",
    cor_titulo=(10, 50, 100) ,
    descricao="Governança de riscos impulsiona \ndecisões mais seguras.")
    tracar_linha_horizontal(slide, y=4.77   , x_start=3.45, comprimento=3.55, texto_circulo= f"{percentuais.get('Informacao', 0)}%", titulo="INFORMAÇÃO",
    cor_titulo=(10, 50, 100) ,
    descricao="Políticas protegem dados \ne garante")
    # === FUNDO BRANCO === #
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

   # === PARALELOGRAMO AZUL  === #
    para_azul = slide.shapes.add_shape(
    MSO_SHAPE.PARALLELOGRAM,
    Inches(0.2),    # Left
    Inches(0.20),   # Top
    Inches(0.5),    # Width
    Inches(0.5)     # Height
)
    para_azul.fill.solid()
    para_azul.fill.fore_color.rgb = RGBColor(30, 115, 190)  # Azul
    para_azul.line.fill.background()

    tf_azul = para_azul.text_frame
    tf_azul.text = "5"  
    p_azul = tf_azul.paragraphs[0]
    p_azul.font.name = "Arial"
    p_azul.font.size = Pt(20)
    p_azul.font.bold = True
    p_azul.font.color.rgb = RGBColor(255, 255, 255)
    p_azul.alignment = PP_ALIGN.CENTER
    tf_azul.vertical_anchor = MSO_ANCHOR.MIDDLE

    # === PARALELOGRAMO CINZA  === #
  
    para_cinza = slide.shapes.add_shape(
    MSO_SHAPE.PARALLELOGRAM,
    Inches(0.85),   # Left
    Inches(0.20),   # Top
    Inches(3.5),    # Width
    Inches(0.5)     # Height
)
    para_cinza.fill.solid()
    para_cinza.fill.fore_color.rgb = RGBColor(70, 70, 70)  
    para_cinza.line.fill.background()

    tf_cinza = para_cinza.text_frame
    tf_cinza.text = "RESUMO EXECUTIVO"  
    tf_cinza.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_cinza.margin_left = Inches(0.2)
    p_cinza = tf_cinza.paragraphs[0]
    p_cinza.alignment = PP_ALIGN.LEFT
    p_cinza.font.name = "Arial"
    p_cinza.font.size = Pt(14)
    p_cinza.font.bold = True
    p_cinza.font.color.rgb = RGBColor(255, 255, 255)

    # === IMAGEM ICON_5EIXOS.png (ocupar toda área disponível) === #
    current_dir = os.path.dirname(__file__)
    parent_dir = os.path.dirname(current_dir)
    icon_path = os.path.join(parent_dir, "images", "5EIXOS.png")
    
    # Área disponível: da direita da caixa até o fim do slide   
    # Largura disponível: começa em 2.7" e vai até quase o fim
    # Altura disponível: abaixo dos paralelogramos até quase o fim
    add_local_image(
        slide,
        icon_path,
        left=1.05,        # ← POSIÇÃO HORIZONTAL
        top=0.7,         # ← POSIÇÃO VERTICAL
        width=5.5,       # ← LARGURA
        height=4.2     # ALTURA
    )

    # === CAIXA ESQUERDA (2.7" largura x 2.1" altura - 3/8 da página) === #
    # Altura da página = 5.625"
    # 3/8 de 5.625 = 2.109375"
    caixa_altura = 1.9
    caixa_largura = 2.2
    
    # Posição: canto inferior esquerdo
    caixa_left = 0.3
    caixa_top = 5.625 - caixa_altura - 0.15  # Bottom padding de 0.3"
    
    caixa = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(caixa_left),
        Inches(caixa_top),
        Inches(caixa_largura),
        Inches(caixa_altura)
    )
    
    # Fundo BRANCO (alterado)
    caixa.fill.solid()
    caixa.fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Borda azul
    caixa.line.fill.solid()
    caixa.line.color.rgb = RGBColor(0, 102, 204)
    caixa.line.width = Pt(2)
    
    # === TEXTO SUPERIOR (azul escuro) + NÚMERO === #
    # Pegar total de perguntas dos dados
   
    # Pegar o array de respostas e contar quantas existem
    respostas = dados.get('dados_modelo', {}).get('respostas', [])
    total_perguntas = len(respostas)

    texto_superior_box = slide.shapes.add_textbox(
        Inches(caixa_left + 0.05),
        Inches(caixa_top + 0.05),
        Inches(caixa_largura - 0.2),
        Inches(0.4)
    )
    tf_sup = texto_superior_box.text_frame
    tf_sup.clear()
    tf_sup.word_wrap = True
    
    p_sup = tf_sup.paragraphs[0]
    p_sup.text = f"Itens não aderentes conforme METODOLOGIA TOTALIZANTE {total_perguntas}"
    p_sup.font.name = "Arial"
    p_sup.font.size = Pt(8)
    p_sup.font.bold = True
    p_sup.font.color.rgb = RGBColor(0, 51, 102) 
    p_sup.alignment = PP_ALIGN.CENTER
    
   # === CAIXA ESQUERDA (2.7" largura x 2.1" altura - 3/8 da página) === #
    # Altura da página = 5.625"
    # 3/8 de 5.625 = 2.109375"
    caixa_altura = 1.9
    caixa_largura = 2.2
    
    # Posição: canto inferior esquerdo
    caixa_left = 0.3
    caixa_top = 5.625 - caixa_altura - 0.15  # Bottom padding de 0.3"
    
    caixa = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(caixa_left),
        Inches(caixa_top),
        Inches(caixa_largura),
        Inches(caixa_altura)
    )
    
    # Fundo BRANCO (alterado)
    caixa.fill.solid()
    caixa.fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Borda azul
    caixa.line.fill.solid()
    caixa.line.color.rgb = RGBColor(0, 102, 204)
    caixa.line.width = Pt(2)
    
    # === TEXTO SUPERIOR (azul escuro) + NÚMERO === #
    # Pegar total de perguntas dos dados
   
    # Pegar o array de respostas e contar quantas existem
    respostas = dados.get('dados_modelo', {}).get('respostas', [])
    total_perguntas = len(respostas)

    texto_superior_box = slide.shapes.add_textbox(
        Inches(caixa_left + 0.05),
        Inches(caixa_top + 0.05),
        Inches(caixa_largura - 0.2),
        Inches(0.4)
    )
    tf_sup = texto_superior_box.text_frame
    tf_sup.clear()
    tf_sup.word_wrap = True
    
    p_sup = tf_sup.paragraphs[0]
    p_sup.text = f"Itens não aderentes conforme METODOLOGIA TOTALIZANTE {total_perguntas}"
    p_sup.font.name = "Arial"
    p_sup.font.size = Pt(8)
    p_sup.font.bold = True
    p_sup.font.color.rgb = RGBColor(0, 51, 102) 
    p_sup.alignment = PP_ALIGN.CENTER
    
    # === RETÂNGULOS AZUIS COM TEXTOS === #
    respostas = dados.get('dados_modelo', {}).get('respostas', [])

    # NC² - Total de não conformidades (vulnerabilidade > 1)
    nao_conformidades = [r for r in respostas if int(r.get('vulnerabilidade', 0)) > 1]
    total_nc = len(nao_conformidades)

    # Recomendações - Total de recomendações
    recomendacoes = [
        r for r in nao_conformidades
        if r.get('recomendacao') and str(r.get('recomendacao')).strip() != ''
    ]
    total_recomendacoes = len(recomendacoes)

    # Posição inicial dos retângulos
    rect_top_start = caixa_top + 0.6
    rect_spacing = 0.42
    rect_height = 0.35

    # Larguras
    label_width = 1.30
    valor_width = caixa_largura - 0.3 - label_width

    itens = [
        ("NC²", total_nc),
        ("Riscos", "Altere"),
        ("Recomendações", total_recomendacoes)
    ]

    for i, (label, valor) in enumerate(itens):
        rect_top = rect_top_start + (i * rect_spacing)

        
        rect_label = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(caixa_left + 0.15),
            Inches(rect_top),
            Inches(label_width),
            Inches(rect_height)
        )
        rect_label.fill.solid()
        rect_label.fill.fore_color.rgb = RGBColor(0, 102, 204)  # Azul claro
        rect_label.line.fill.background()

        
        rect_valor = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(caixa_left + 0.15 + label_width),
            Inches(rect_top),
            Inches(valor_width),
            Inches(rect_height)
        )
        rect_valor.fill.solid()
        rect_valor.fill.fore_color.rgb = RGBColor(10, 50, 100)  # Azul escuro
        rect_valor.line.fill.background()

        #  TEXTBOX DO LABEL
        label_box = slide.shapes.add_textbox(
            Inches(caixa_left + 0.25),
            Inches(rect_top),
            Inches(label_width - 0.1),
            Inches(rect_height)
        )
        tf_label = label_box.text_frame
        tf_label.clear()
        tf_label.word_wrap = False
        tf_label.vertical_anchor = MSO_ANCHOR.MIDDLE

        p_label = tf_label.paragraphs[0]
        p_label.text = label
        p_label.font.name = "Arial"
        p_label.font.size = Pt(9)
        p_label.font.bold = True
        p_label.font.color.rgb = RGBColor(255, 255, 255)
        p_label.alignment = PP_ALIGN.LEFT

        # TEXTBOX DO VALOR
        valor_box = slide.shapes.add_textbox(
            Inches(caixa_left + 0.15 + label_width + 0.05),
            Inches(rect_top),
            Inches(valor_width - 0.1),
            Inches(rect_height)
        )
        tf_valor = valor_box.text_frame
        tf_valor.clear()
        tf_valor.word_wrap = False
        tf_valor.vertical_anchor = MSO_ANCHOR.MIDDLE

        p_valor = tf_valor.paragraphs[0]
        p_valor.text = str(valor)
        p_valor.font.name = "Arial"
        p_valor.font.size = Pt(9)
        p_valor.font.bold = True
        p_valor.font.color.rgb = RGBColor(255, 255, 255)
        p_valor.alignment = PP_ALIGN.RIGHT
    
    valores = [
    percentuais.get('Tecnologia', 0),
    percentuais.get('Processos', 0),
    percentuais.get('Pessoas', 0),
    percentuais.get('Gestao', 0),
    percentuais.get('Informacao', 0),
]

    media_percentual = round(sum(valores) / len(valores))

    texto_70 = slide.shapes.add_textbox(
    Inches(2.8),  # x
    Inches(2.55),  # y
    Inches(1.0),
    Inches(0.6)
)

    tf_70 = texto_70.text_frame
    tf_70.clear()
    tf_70.vertical_anchor = MSO_ANCHOR.MIDDLE

    p_70 = tf_70.paragraphs[0]
    p_70.text = f"{media_percentual}%"
    p_70.font.size = Pt(36)
    p_70.font.bold = True
    p_70.font.color.rgb = RGBColor(255, 255, 255)
    p_70.alignment = PP_ALIGN.CENTER

    print("✅ Slide de Resumo Executivo criado!", file=sys.stderr)
    return slide