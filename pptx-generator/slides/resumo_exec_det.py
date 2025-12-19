#!/usr/bin/env python3
# -*- coding: utf-8 -*-
import sys
import os
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


def add_local_image(slide, image_path, left, top, width, height):
    """Adiciona imagem local ao slide"""
    try:
        if os.path.exists(image_path):
            slide.shapes.add_picture(
                image_path,
                Inches(left),
                Inches(top),
                width=Inches(width),
                height=Inches(height)
            )
            return True
        else:
            print(f"⚠️ Imagem não encontrada: {image_path}", file=sys.stderr)
            return False
    except Exception as e:
        print(f"❌ Erro ao adicionar imagem: {e}", file=sys.stderr)
        return False

def gerar_resumo_exec_det(pres, dados):

    slide = pres.slides.add_slide(pres.slide_layouts[6])
         # === FUNDO BRANCO === #
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

   # === PARALELOGRAMO AZUL  === #
    para_azul = slide.shapes.add_shape(
    MSO_SHAPE.PARALLELOGRAM,
    Inches(0.2),    # Left
    Inches(0.20),   # Top
    Inches(0.5),    # Width
    Inches(0.5)     # Height
)
    para_azul.fill.solid()
    para_azul.fill.fore_color.rgb = RGBColor(30, 115, 190)  # Azul
    para_azul.line.fill.background()

    tf_azul = para_azul.text_frame
    tf_azul.text = "6"  
    p_azul = tf_azul.paragraphs[0]
    p_azul.font.name = "Arial"
    p_azul.font.size = Pt(20)
    p_azul.font.bold = True
    p_azul.font.color.rgb = RGBColor(255, 255, 255)
    p_azul.alignment = PP_ALIGN.CENTER
    tf_azul.vertical_anchor = MSO_ANCHOR.MIDDLE

    # === PARALELOGRAMO CINZA  === #
  
    para_cinza = slide.shapes.add_shape(
    MSO_SHAPE.PARALLELOGRAM,
    Inches(0.85),   # Left
    Inches(0.20),   # Top
    Inches(5.0),    # Width
    Inches(0.5)     # Height
)
    para_cinza.fill.solid()
    para_cinza.fill.fore_color.rgb = RGBColor(70, 70, 70)  
    para_cinza.line.fill.background()

    tf_cinza = para_cinza.text_frame
    tf_cinza.text = "RESUMO EXECUTIVO - DETALHADO"  
    tf_cinza.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_cinza.margin_left = Inches(0.2)
    p_cinza = tf_cinza.paragraphs[0]
    p_cinza.alignment = PP_ALIGN.LEFT
    p_cinza.font.name = "Arial"
    p_cinza.font.size = Pt(14)
    p_cinza.font.bold = True
    p_cinza.font.color.rgb = RGBColor(255, 255, 255)


    # === SUBTÍTULO === #
    subtitulo = slide.shapes.add_textbox(
        Inches(2.5),
        Inches(0.75),
        Inches(5.0),
        Inches(0.35)
    )

    tf_sub = subtitulo.text_frame
    tf_sub.clear()
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "Itens Aderentes Conforme Metodologia"
    p_sub.font.name = "Arial"
    p_sub.font.size = Pt(15)
    p_sub.font.bold = True
    p_sub.font.color.rgb = RGBColor(0, 38, 77)
    p_sub.alignment = PP_ALIGN.CENTER

    # === RETÂNGULOS VERTICAIS === #
    altura_ret = 4.0
    largura_ret = 1.7
    margem_lateral = 0.45
    espaco = 0.15
    topo = 1.45  # abaixo do subtítulo

    cores = [
        RGBColor(0, 51, 102),   
        RGBColor(0, 76, 153),    
        RGBColor(35, 137, 239),    
        RGBColor(102, 178, 255),    
        RGBColor(0, 180, 170)    
    ]

    titulos = [
    "TECNOLOGIA",
    "PROCESSOS",
    "PESSOAL",
    "GESTÃO",
    "INFORMAÇÃO"
    ]
    chaves_pilares = ["Tecnologia", "Processos", "Pessoas", "Gestao", "Informacao"]

    current_dir = os.path.dirname(__file__)      # pptx-generator/slides/
    parent_dir = os.path.dirname(current_dir)    # pptx-generator/

    icons = [
        os.path.join(parent_dir, "images", "ICON_TECNOLOGIA.png"),
        os.path.join(parent_dir, "images", "ICON_PROCESSOS.png"),
        os.path.join(parent_dir, "images", "ICON_PESSOAS.png"),
        os.path.join(parent_dir, "images", "ICON_GESTAO.png"),
        os.path.join(parent_dir, "images", "ICON_INFORMACOES.png"),
    ]

    for i in range(5):

        left = margem_lateral + i * (largura_ret + espaco)

        raio = 0.35
        diametro = raio * 2
        topo_circulo = 1.1
        altura_titulo = 0.6
        margem_vertical = 0.15

        pilar = chaves_pilares[i]
        topicos_pilar = dados.get("dados_modelo", {}).get("analise_topicos", {}).get(pilar, [])

        # === RETÂNGULO === #
        ret = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(topo),
            Inches(largura_ret),
            Inches(altura_ret)
        )

        ret.fill.solid()
        ret.fill.fore_color.rgb = cores[i]
        ret.line.fill.background()

        # === CÍRCULO BRANCO === #
        raio = 0.35
        diametro = raio * 2

        centro_x = left + (largura_ret / 2)
        topo_circulo = 1.1

        circulo = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(centro_x - raio),
            Inches(topo_circulo),
            Inches(diametro),
            Inches(diametro)
        )

        circulo.fill.solid()
        circulo.fill.fore_color.rgb = cores[i]
        circulo.line.fill.solid()
        circulo.line.fill.fore_color.rgb = RGBColor(255, 255, 255)
        circulo.line.width = Pt(2)

        icone_tamanho = 0.6  # tamanho do ícone     

        icone_left = left + (largura_ret - icone_tamanho) / 2
        icone_top = topo - 0.3

        if os.path.exists(icons[i]):
            slide.shapes.add_picture(
                icons[i],
                Inches(icone_left),
                Inches(icone_top),
                width=Inches(icone_tamanho),
                height=Inches(icone_tamanho)
            )
        else:
            print(f"⚠️ Ícone não encontrado: {icons[i]}", file=sys.stderr)

        # === TÍTULO DO RETÂNGULO === #
        titulo_box = slide.shapes.add_textbox(
            Inches(left),
            Inches(topo_circulo + diametro + 0.15),
            Inches(largura_ret),
            Inches(0.6)
        )

        tf_titulo = titulo_box.text_frame
        tf_titulo.clear()
        tf_titulo.vertical_anchor = MSO_ANCHOR.TOP

        p_titulo = tf_titulo.paragraphs[0]
        p_titulo.text = titulos[i]
        p_titulo.font.name = "Arial"
        p_titulo.font.size = Pt(12)
        p_titulo.font.bold = True
        p_titulo.font.color.rgb = RGBColor(255, 255, 255)
        p_titulo.alignment = PP_ALIGN.CENTER

        margem_interna = 0.15
        qtd_ret_internos = 4
        altura_ret_interno = 0.55
        espaco_vertical = 0.15

        largura_interna = largura_ret - (margem_interna * 2)

        titulo_top = topo_circulo + diametro + margem_vertical

        topo_inicial_interno = (
            titulo_top +
            altura_titulo +
            margem_vertical
        )

        for j in range(qtd_ret_internos):
            top_interno = topo_inicial_interno + j * (altura_ret_interno + espaco_vertical) - 0.3

            ret_interno = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left + margem_interna),
                Inches(top_interno),
                Inches(largura_interna),
                Inches(altura_ret_interno)
            )

            ret_interno.fill.solid()
            ret_interno.fill.fore_color.rgb = RGBColor(255, 255, 255)
            ret_interno.line.fill.background()

            # Se houver tópico para essa posição, adiciona o texto
            if j < len(topicos_pilar):
                topico = topicos_pilar[j]
                nome_topico = topico.get("topico_nome", "")
                fracao = topico.get("fracao", "0/0")
               
                #nome_topico = nome_topico.encode('latin1').decode('utf-8')
                tf = ret_interno.text_frame
                tf.clear()
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE

                p = tf.paragraphs[0]
                p.text = f"{nome_topico}\n{fracao}"
                p.alignment = PP_ALIGN.CENTER

                p.font.name = "Arial"
                p.font.size = Pt(11)
                p.font.bold = True
                p.font.color.rgb = cores[i]

            else:
                tf = ret_interno.text_frame
                tf.clear()
    print("✅ Slide de Resumo Executivo Detalhado criado!", file=sys.stderr)
    return slide 

