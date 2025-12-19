import sys
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL
from pptx.oxml.xmlchemy import OxmlElement


def add_base64_image(slide, base64_str, left, top, width, height):
    """Adiciona imagem base64 ao slide"""
    try:
        from io import BytesIO
        import base64
        
        if ',' in base64_str:
            base64_str = base64_str.split(',')[1]

        img_data = base64.b64decode(base64_str)
        stream = BytesIO(img_data)
        slide.shapes.add_picture(
            stream, Inches(left), Inches(top),
            width=Inches(width), height=Inches(height)
        )
        return True
    except Exception as e:
        print(f"Erro ao adicionar imagem: {e}", file=sys.stderr)
        return False


def create_diagonal_polygon(slide, pres):
    """Cria polígono diagonal azul usando coordenadas customizadas"""
    # Converter para EMUs (English Metric Units - unidade do PowerPoint)
    width = pres.slide_width
    height = pres.slide_height
    
    # Criar o shape como freeform
    shapes = slide.shapes
    shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(1), Inches(1))
    
    # Ajustar para polígono customizado via XML
    sp = shape._element
    spPr = sp.spPr
    
    # Remover geometria padrão
    for child in list(spPr):
        if 'prstGeom' in child.tag:
            spPr.remove(child)
    
    # Criar geometria customizada
    custGeom = OxmlElement('a:custGeom')
    
    # Adicionar avLst (lista de variáveis de ajuste - vazio para polígono simples)
    avLst = OxmlElement('a:avLst')
    custGeom.append(avLst)
    
    # Adicionar gdLst (lista de guias - vazio)
    gdLst = OxmlElement('a:gdLst')
    custGeom.append(gdLst)
    
    # Adicionar ahLst (lista de alças - vazio)
    ahLst = OxmlElement('a:ahLst')
    custGeom.append(ahLst)
    
    # Adicionar cxnLst (lista de conexões - vazio)
    cxnLst = OxmlElement('a:cxnLst')
    custGeom.append(cxnLst)
    
    # Adicionar pathLst (lista de caminhos)
    pathLst = OxmlElement('a:pathLst')
    path = OxmlElement('a:path')
    path.set('w', str(width))
    path.set('h', str(height))
    
    # Adicionar vértices
    # MoveTo (primeiro ponto)
    moveTo = OxmlElement('a:moveTo')
    pt = OxmlElement('a:pt')
    pt.set('x', '0')
    pt.set('y', '0')
    moveTo.append(pt)
    path.append(moveTo)
    
    # LineTo para os outros pontos
    # Criando forma diagonal: topo direito mais largo, base direita mais estreita
    points = [
        (int(Inches(5.5)), 0),           # Topo direito da diagonal (mais largo)
        (int(Inches(2.5)), int(height)), # Base direita da diagonal (mais estreita - mais inclinado)
        (0, int(height))                 # Base esquerda
    ]
    
    for x, y in points:
        lineTo = OxmlElement('a:lnTo')
        pt = OxmlElement('a:pt')
        pt.set('x', str(x))
        pt.set('y', str(y))
        lineTo.append(pt)
        path.append(lineTo)
    
    # Fechar o caminho
    close = OxmlElement('a:close')
    path.append(close)
    
    pathLst.append(path)
    custGeom.append(pathLst)
    
    # Adicionar ao shape
    spPr.append(custGeom)
    
    # Definir posição e tamanho
    shape.left = 0
    shape.top = 0
    shape.width = width
    shape.height = height
    
    # Aplicar preenchimento com gradiente azul
    fill = shape.fill
    fill.gradient()
    fill.gradient_angle = 90
    stop_dark = fill.gradient_stops[0]
    stop_dark.position = 0.0
    stop_dark.color.rgb = RGBColor(10, 50, 100)
    stop_light = fill.gradient_stops[1]
    stop_light.position = 1.0
    stop_light.color.rgb = RGBColor(30, 115, 190)
    
    # Remover linha
    shape.line.fill.background()
    
    return shape


def create_gray_triangle(slide, pres):
    """Cria triângulo cinza isósceles na parte inferior (1/3 da página)"""
    width = pres.slide_width
    height = pres.slide_height
    
    # Calcular 1/3 da altura do slide
    triangle_height = height / 3
    triangle_start_y = height - triangle_height  # Começa a 2/3 da página
    
    # Calcular onde está a borda da faixa azul nessa altura
    # A faixa azul vai de 5.5 inches no topo até 2.5 inches na base
    # Interpolação linear para encontrar a posição X na altura do vértice
    top_x = Inches(5.5)
    bottom_x = Inches(2.5)
    
    # Proporção da altura onde o triângulo começa (2/3 = 0.666)
    ratio = triangle_start_y / height
    diagonal_x = top_x + (bottom_x - top_x) * ratio
    
    # Largura da base do triângulo isósceles (ajustar para formar isósceles)
    base_left_x = bottom_x  # Onde a faixa azul termina na base
    base_width = Inches(3.0)  # Largura da base do triângulo
    base_right_x = base_left_x + base_width
    
    # Criar o shape
    shapes = slide.shapes
    shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(1), Inches(1))
    
    # Ajustar para polígono customizado via XML
    sp = shape._element
    spPr = sp.spPr
    
    # Remover geometria padrão
    for child in list(spPr):
        if 'prstGeom' in child.tag:
            spPr.remove(child)
    
    # Criar geometria customizada
    custGeom = OxmlElement('a:custGeom')
    
    # Adicionar listas vazias necessárias
    avLst = OxmlElement('a:avLst')
    custGeom.append(avLst)
    
    gdLst = OxmlElement('a:gdLst')
    custGeom.append(gdLst)
    
    ahLst = OxmlElement('a:ahLst')
    custGeom.append(ahLst)
    
    cxnLst = OxmlElement('a:cxnLst')
    custGeom.append(cxnLst)
    
    # Adicionar pathLst (lista de caminhos)
    pathLst = OxmlElement('a:pathLst')
    path = OxmlElement('a:path')
    path.set('w', str(width))
    path.set('h', str(height))
    
    # Vértices do triângulo ISÓSCELES:
    # Ponto 1: Vértice superior (na linha da diagonal, a 2/3 da altura)
    # Ponto 2: Vértice inferior esquerdo (onde termina a faixa azul na base)
    # Ponto 3: Vértice inferior direito (formando base do triângulo)
    
    # MoveTo (vértice superior - ponta do triângulo)
    moveTo = OxmlElement('a:moveTo')
    pt = OxmlElement('a:pt')
    pt.set('x', str(int(diagonal_x)))
    pt.set('y', str(int(triangle_start_y)))
    moveTo.append(pt)
    path.append(moveTo)
    
    # LineTo para os outros pontos
    points = [
        (int(base_left_x), int(height)),   # Base esquerda
        (int(base_right_x), int(height))   # Base direita
    ]
    
    for x, y in points:
        lineTo = OxmlElement('a:lnTo')
        pt = OxmlElement('a:pt')
        pt.set('x', str(x))
        pt.set('y', str(y))
        lineTo.append(pt)
        path.append(lineTo)
    
    # Fechar o caminho
    close = OxmlElement('a:close')
    path.append(close)
    
    pathLst.append(path)
    custGeom.append(pathLst)
    
    # Adicionar ao shape
    spPr.append(custGeom)
    
    # Definir posição e tamanho
    shape.left = 0
    shape.top = 0
    shape.width = width
    shape.height = height
    
    # Aplicar preenchimento cinza escuro
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(70, 70, 70)
    
    # Remover linha
    shape.line.fill.background()
    
    return shape


def create_white_triangle(slide, pres):
    """Cria triângulo retângulo branco na quina superior esquerda (2/3 altura x 2.5 inches base)"""
    width = pres.slide_width
    height = pres.slide_height
    
    # Calcular 2/3 da altura do slide
    triangle_height = height * 4 / 5
    
    # Base do triângulo
    triangle_base = Inches(2.5)
    
    # Criar o shape
    shapes = slide.shapes
    shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(1), Inches(1))
    
    # Ajustar para polígono customizado via XML
    sp = shape._element
    spPr = sp.spPr
    
    # Remover geometria padrão
    for child in list(spPr):
        if 'prstGeom' in child.tag:
            spPr.remove(child)
    
    # Criar geometria customizada
    custGeom = OxmlElement('a:custGeom')
    
    # Adicionar listas vazias necessárias
    avLst = OxmlElement('a:avLst')
    custGeom.append(avLst)
    
    gdLst = OxmlElement('a:gdLst')
    custGeom.append(gdLst)
    
    ahLst = OxmlElement('a:ahLst')
    custGeom.append(ahLst)
    
    cxnLst = OxmlElement('a:cxnLst')
    custGeom.append(cxnLst)
    
    # Adicionar pathLst (lista de caminhos)
    pathLst = OxmlElement('a:pathLst')
    path = OxmlElement('a:path')
    path.set('w', str(width))
    path.set('h', str(height))
    
    # Vértices do triângulo RETÂNGULO:
    # Ponto 1: Quina superior esquerda (0, 0)
    # Ponto 2: Base do triângulo no topo (2.5 inches, 0)
    # Ponto 3: Vértice inferior na borda esquerda (0, 2/3 altura)
    
    # MoveTo (quina superior esquerda)
    moveTo = OxmlElement('a:moveTo')
    pt = OxmlElement('a:pt')
    pt.set('x', '0')
    pt.set('y', '0')
    moveTo.append(pt)
    path.append(moveTo)
    
    # LineTo para os outros pontos
    points = [
        (int(triangle_base), 0),           # Base no topo (direita)
        (0, int(triangle_height))          # Vértice inferior (esquerda)
    ]
    
    for x, y in points:
        lineTo = OxmlElement('a:lnTo')
        pt = OxmlElement('a:pt')
        pt.set('x', str(x))
        pt.set('y', str(y))
        lineTo.append(pt)
        path.append(lineTo)
    
    # Fechar o caminho
    close = OxmlElement('a:close')
    path.append(close)
    
    pathLst.append(path)
    custGeom.append(pathLst)
    
    # Adicionar ao shape
    spPr.append(custGeom)
    
    # Definir posição e tamanho
    shape.left = 0
    shape.top = 0
    shape.width = width
    shape.height = height
    
    # Aplicar preenchimento branco
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Remover linha
    shape.line.fill.background()
    
    return shape


def gerar_sumario(pres, dados):
    """Gera o slide de sumário"""
    print("Criando slide de sumário...", file=sys.stderr)
    slide = pres.slides.add_slide(pres.slide_layouts[6])

    # === FUNDO BRANCO === #
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # === ÁREA AZUL DIAGONAL (usando polígono customizado) === #
    create_diagonal_polygon(slide, pres)
    
    # === TRIÂNGULO CINZA ISÓSCELES (rente à faixa azul) === #
    create_gray_triangle(slide, pres)
    
    # === TRIÂNGULO RETÂNGULO BRANCO (quina superior esquerda) === #
    create_white_triangle(slide, pres)

    # === LOGO POWER (topo esquerdo sobre o fundo branco) === #
    logo = dados.get("imagens", {}).get("logo_empresa")
    if logo:
        add_base64_image(slide, logo, 0.3, 0.3, 1.5, 0.9)

    # === TÍTULO "SUMÁRIO" (centralizado na faixa azul) === #
    titulo_box = slide.shapes.add_textbox(
        Inches(1.0),
        Inches(2.5),
        Inches(3.0),
        Inches(0.6)
    )
    tf = titulo_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "SUMÁRIO"
    p.font.name = "Arial"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # === LISTA DE TÓPICOS === #
    topicos = [
        "OBJETIVO",
        "METODOLOGIA",
        "PANORAMA SITUACIONAL",
        "SEGURANÇA PÚBLICA",
        "RESUMO EXECUTIVO",
        "ITENS ADERENTES",
        "VULNERABILIDADES",
        "RECOMENDAÇÕES"
    ]

    start_left = Inches(5.6)
    start_top = Inches(0.5)
    spacing = Inches(0.55)

    for i, topico in enumerate(topicos):
        top_pos = start_top + (i * spacing)
        
        # Coordenadas principais do paralelogramo
        para_left = start_left
        para_top = top_pos
        para_width = Inches(0.5)
        para_height = Inches(0.35)

        # === "Sombra" manual: paralelogramo cinza atrás ===
        shadow_para = slide.shapes.add_shape(
            MSO_SHAPE.PARALLELOGRAM,
            para_left + Inches(0.04),
            para_top + Inches(0.04),
            para_width,
            para_height
        )
        shadow_para.fill.solid()
        shadow_para.fill.fore_color.rgb = RGBColor(90,90,90)
        shadow_para.line.fill.background()

        # === Paralelogramo azul principal ===
        para = slide.shapes.add_shape(
            MSO_SHAPE.PARALLELOGRAM,
            para_left,
            para_top,
            para_width,
            para_height
        )
        para.fill.solid()
        para.fill.fore_color.rgb = RGBColor(30, 115, 190)
        para.line.fill.background()

        # Número dentro do paralelogramo azul
        tf_para = para.text_frame
        tf_para.clear()
        tf_para.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_para = tf_para.paragraphs[0]
        p_para.text = str(i + 1)
        p_para.font.name = "Calibri"
        p_para.font.size = Pt(16)
        p_para.font.bold = True
        p_para.font.color.rgb = RGBColor(255, 255, 255)
        p_para.alignment = PP_ALIGN.CENTER

        # === CAIXA BRANCA COM TEXTO === #
        box_left = para_left + para_width + Inches(0.1)
        box_top = top_pos
        box_width = Inches(3.5)
        box_height = Inches(0.35)
        
        texto_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            box_left,
            box_top,
            box_width,
            box_height
        )
        texto_box.fill.solid()
        texto_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        texto_box.line.color.rgb = RGBColor(200, 200, 200)
        texto_box.line.width = Pt(1)
        
        tf_box = texto_box.text_frame
        tf_box.clear()
        tf_box.margin_left = Inches(0.15)
        tf_box.margin_top = Inches(0.03)
        p_box = tf_box.paragraphs[0]
        p_box.text = topico
        p_box.font.name = "Calibri"
        p_box.font.size = Pt(14)
        p_box.font.bold = True
        p_box.font.color.rgb = RGBColor(30, 115, 190)
        p_box.alignment = PP_ALIGN.LEFT
        tf_box.vertical_anchor = MSO_ANCHOR.MIDDLE

    print("Slide de sumário criado!", file=sys.stderr)
    return slide