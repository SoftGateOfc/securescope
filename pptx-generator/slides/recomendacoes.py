#!/usr/bin/env python3
# -*- coding: utf-8 -*-
import sys
import os
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import base64
from io import BytesIO

# Importar funções de processamento
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..'))
from utils.processar_dados import reorganizar_por_prioridade, mapear_icone_pilar, mapear_cor_prioridade  # ← ADICIONAR nomear_nivel_vulnerabilidade se precisar


def add_local_image(slide, image_path, left, top, width, height):
    """Adiciona imagem local ao slide"""
    try:
        if os.path.exists(image_path):
            slide.shapes.add_picture(
                image_path, 
                Inches(left), 
                Inches(top),
                width=Inches(width), 
                height=Inches(height)
            )
            return True
        else:
            print(f"Imagem não encontrada: {image_path}", file=sys.stderr, flush=True)
            return False
    except Exception as e:
        print(f"Erro ao adicionar imagem local: {e}", file=sys.stderr, flush=True)
        return False


def criar_slide_recomendacoes(pres, dados, itens_slide):

    # Altura por linha: 0.85 se recomendação > 120 chars, senão 0.64
    alturas_linhas = [0.85 if len(item['recomendacao']) > 120 else 0.64 for item in itens_slide]
    altura_total_tabela = 0.15 + sum(alturas_linhas)

    slide = pres.slides.add_slide(pres.slide_layouts[6])

    # === FUNDO BRANCO === #
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # === PARALELOGRAMO AZUL === #
    para_azul = slide.shapes.add_shape(
        MSO_SHAPE.PARALLELOGRAM,
        Inches(0.2), Inches(0.2), Inches(0.5), Inches(0.5)
    )
    para_azul.fill.solid()
    para_azul.fill.fore_color.rgb = RGBColor(30, 115, 190)
    para_azul.line.fill.background()

    tf_azul = para_azul.text_frame
    tf_azul.text = "9"  
    p = tf_azul.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    tf_azul.vertical_anchor = MSO_ANCHOR.MIDDLE

    # === PARALELOGRAMO CINZA === #
    para_cinza = slide.shapes.add_shape(
        MSO_SHAPE.PARALLELOGRAM,
        Inches(0.85), Inches(0.2), Inches(5.0), Inches(0.5)
    )
    para_cinza.fill.solid()
    para_cinza.fill.fore_color.rgb = RGBColor(70, 70, 70)
    para_cinza.line.fill.background()

    tf_cinza = para_cinza.text_frame
    tf_cinza.text = "RECOMENDAÇÕES"
    tf_cinza.margin_left = Inches(0.2)
    tf_cinza.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_cinza.paragraphs[0].font.bold = True
    tf_cinza.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # === SUBTÍTULO === #
    subtitulo = slide.shapes.add_textbox(
        Inches(2.5), Inches(0.75), Inches(5.0), Inches(0.35)
    )
    p = subtitulo.text_frame.paragraphs[0]
    p.text = "Prioridades de Melhorias"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 38, 77)
    p.alignment = PP_ALIGN.CENTER

    # ================= TABELA ================= #
    rows = 1 + len(itens_slide)
    cols = 6  

    table = slide.shapes.add_table(
        rows, cols,
        Inches(0.25), Inches(1.2),
        Inches(9.05),
        Inches(altura_total_tabela)
    ).table

    # Larguras (6 colunas)
    widths = [0.9, 0.5, 2.05, 2.6 , 2.1, 0.9]
    for i, w in enumerate(widths):
        table.columns[i].width = Inches(w)

    # Alturas
    table.rows[0].height = Inches(0.3 / 2)
    for i in range(1, rows):
        table.rows[i].height = Inches(alturas_linhas[i - 1])

    # Cabeçalho
    headers = [
        "Pilares", "NC", "Classificação",
        "Recomendação", "Riscos Mitigados", "Prioridade"  
    ]

    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(30, 115, 190)

        p = cell.text_frame.paragraphs[0]
        p.text = header
        p.font.bold = True
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # ================= RENDERIZAR LINHAS DA TABELA ================= #
    # Caminho das imagens
    current_dir = os.path.dirname(__file__)
    parent_dir = os.path.dirname(current_dir)
    images_dir = os.path.join(parent_dir, "images")
    
    # Guardar posição da tabela para calcular posições das células
    table_left = Inches(0.3)
    table_top = Inches(1.2)
    header_height = Inches(0.3)

    for row_idx, item in enumerate(itens_slide, start=1):
        # Calcular posição Y desta linha somando alturas anteriores
        cell_top = table_top + header_height + sum(Inches(alturas_linhas[j]) for j in range(row_idx - 1))

        # Alternar cores das linhas
        if row_idx % 2 == 0:
            cor_linha = RGBColor(230, 240, 250) # azul
        else:
            cor_linha = RGBColor(220, 220, 220)  # cinza

        for col in range(cols):
            cell = table.cell(row_idx, col)
            cell.fill.solid()
            cell.fill.fore_color.rgb = cor_linha
        
        # Coluna 0: ÍCONE DO ELEMENTO (pilar)
        cell = table.cell(row_idx, 0)
        
        # Carregar o ícone
        icone_nome = item["elemento"]
        image_path = os.path.join(images_dir, icone_nome)
        
        # Calcular posição X da coluna 0
        col_0_left = table_left
        
        try:
            if os.path.exists(image_path):
                slide.shapes.add_picture(
                    image_path,
                    col_0_left + Inches(0.15),
                    cell_top + Inches(0.1),  
                    width=Inches(0.4)
                )
        except Exception as e:
            print(f"Erro ao carregar ícone {icone_nome}: {e}", file=sys.stderr)
        
        # Coluna 1: NC Atendida
        cell = table.cell(row_idx, 1)
        p = cell.text_frame.paragraphs[0]
        p.text = item["nc"]
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(30, 115, 190)
        p.alignment = PP_ALIGN.CENTER
        cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Coluna 2: Classificação
        cell = table.cell(row_idx, 2)
        p = cell.text_frame.paragraphs[0]
        p.text = item["classificacao"]
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(30, 115, 190)
        p.alignment = PP_ALIGN.LEFT
        cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.text_frame.word_wrap = True
        
        # Coluna 3: Recomendação 
        cell = table.cell(row_idx, 3)
        p = cell.text_frame.paragraphs[0]
        p.text = item["recomendacao"]  
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(30, 115, 190)
        p.alignment = PP_ALIGN.LEFT
        cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.text_frame.word_wrap = True
        
        # Coluna 4: Riscos Mitigados 
        cell = table.cell(row_idx, 4)
        p = cell.text_frame.paragraphs[0]
        p.text = item["riscos_mitigados"]
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Coluna 5: Prioridade (círculo colorido)
        cell = table.cell(row_idx, 5)
        
        # Adicionar círculo colorido
        prioridade = item["prioridade"]
        cor_rgb = mapear_cor_prioridade(prioridade)
        
        # Calcular posição X da coluna 5 (última coluna)
        
        col_5_left = table_left + Inches(8.05)
        
        circulo = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            col_5_left + Inches(0.35),
            cell_top + Inches(0.17), 
            Inches(0.25),
            Inches(0.25)
        )
        circulo.fill.solid()
        circulo.fill.fore_color.rgb = RGBColor(*cor_rgb)
        circulo.line.fill.background()

    
    # -------- TÍTULO DA LEGENDA  -------- #
    titulo_legenda = slide.shapes.add_textbox(
        Inches(0.3),
        Inches(4.75),
        Inches(4.2),
        Inches(0.25)
    )

    p = titulo_legenda.text_frame.paragraphs[0]
    p.text = "Prioridade"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.LEFT

    legenda = slide.shapes.add_textbox(
        Inches(0.2),
        Inches(5.0),
        Inches(4.2),
        Inches(0.4)
    )

    # Borda azul escura
    legenda.line.color.rgb = RGBColor(0, 51, 102)
    legenda.line.width = Pt(0.9)

    tf_legenda = legenda.text_frame
    tf_legenda.clear()
    tf_legenda.margin_left = Inches(0.15)
    tf_legenda.margin_top = Inches(0.05)

    # -------- ITENS DA LEGENDA (HORIZONTAIS) -------- #
    legenda_itens = [
        ("Longo Prazo",  RGBColor(192, 0, 0)),
        ("Médio Prazo",  RGBColor(255, 192, 0)),
        ("Curto Prazo", RGBColor(146, 208, 80)),
    ]

    start_x = Inches(0.3)
    y = Inches(5.1)
    espacamento = Inches(1.2)

    for texto, cor in legenda_itens:
        # Círculo
        circulo = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            start_x, y,
            Inches(0.18), Inches(0.18)
        )
        circulo.fill.solid()
        circulo.fill.fore_color.rgb = cor
        circulo.line.fill.background()

        # Texto ao lado do círculo
        label = slide.shapes.add_textbox(
            start_x + Inches(0.20),
            y - Inches(0.02),
            Inches(0.9),
            Inches(0.25)
        )
        p = label.text_frame.paragraphs[0]
        p.text = texto
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(30, 115, 190)
        p.alignment = PP_ALIGN.LEFT

        start_x += espacamento
      
    titulo_elementos = slide.shapes.add_textbox(
        Inches(4.7),
        Inches(4.75),
        Inches(4.8),
        Inches(0.25)
    )

    p = titulo_elementos.text_frame.paragraphs[0]
    p.text = "Pilares"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.LEFT

    # ========= CAIXA DA LEGENDA ========= #
    legenda_elementos = slide.shapes.add_textbox(
        Inches(4.5),
        Inches(5.0),
        Inches(4.6),
        Inches(0.4)
    )

    # Fundo azul escuro
    legenda_elementos.fill.solid()
    legenda_elementos.fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Borda
    legenda_elementos.line.color.rgb = RGBColor(0, 51, 102)
    legenda_elementos.line.width = Pt(0.9)

    # ========= ITENS DA LEGENDA ========= #
    elementos = [
        ("tecnologia.png",   "Tecnologia"),
        ("processos.png",    "Processos"),
        ("pessoas.png",      "Pessoas"),
        ("informacoes.png",  "Informação"),
        ("gestao.png",       "Gestão"),
    ]

    start_x = Inches(4.6)
    y = Inches(5.05)
    espacamento = Inches(0.95)

    for nome_icone, texto in elementos:
        image_path = os.path.join(images_dir, nome_icone)

        # Ícone
        try:
            slide.shapes.add_picture(
                image_path,
                start_x,
                y,
                width=Inches(0.28)
            )
        except Exception as e:
            print(f"Erro ao carregar ícone {nome_icone}: {e}", file=sys.stderr)

        # Texto ao lado do ícone
        label = slide.shapes.add_textbox(
            start_x + Inches(0.22),
            y - Inches(0.02),
            Inches(0.8),
            Inches(0.3)
        )

        p = label.text_frame.paragraphs[0]
        p.text = texto
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(30, 115, 190)
        p.alignment = PP_ALIGN.LEFT

        start_x += espacamento
    
    # === LOGO === #
    logo = dados.get("imagens", {}).get("logo_empresa")
    if logo:
        try:
            if ',' in logo:
                logo = logo.split(',')[1]
            stream = BytesIO(base64.b64decode(logo))
            slide.shapes.add_picture(stream, Inches(9.15), Inches(4.7), width=Inches(0.8))
        except Exception as e:
            print(f"Erro ao adicionar logo: {e}", file=sys.stderr)

    return slide


def gerar_recomendacoes(pres, dados):
    """
    Gera slides de recomendações com dados reais
    """
    print("🔄 Processando recomendações...", file=sys.stderr, flush=True)
    
    # Puxar dados reais das respostas
    respostas = dados.get('dados_modelo', {}).get('respostas', [])
    
    if not respostas or len(respostas) == 0:
        print("⚠️ Nenhuma resposta encontrada", file=sys.stderr, flush=True)
        return
    
    # Reorganizar por PRIORIDADE (não criticidade)
    recomendacoes = reorganizar_por_prioridade(respostas)
    
    if len(recomendacoes) == 0:
        print("⚠️ Nenhuma recomendação encontrada (todas com nível 1)", file=sys.stderr, flush=True)
        return
    
    print(f"📊 {len(recomendacoes)} recomendações encontradas", file=sys.stderr, flush=True)
    
    # Converter para o formato esperado pela tabela
    itens = []
    for rec in recomendacoes:
        # Extrair nome do pilar para pegar o ícone correto
        pilar_nome = rec.get('pilar', '')
        icone_nome = mapear_icone_pilar(pilar_nome)
        
        # Montar texto da classificação
        nome_topico = rec.get('topicos', 'Sem descrição')
        #   nome_topico = nome_topico.encode('latin1').decode('utf-8')
        classificacao_texto = nome_topico
        
        # Pegar recomendacao e garantir que seja string
        recomendacao_raw = rec.get('recomendacao', '')
        if recomendacao_raw is None or str(recomendacao_raw).strip() == "":
            recomendacao_texto = "Não há recomendações para essa pergunta"
        else:
            recomendacao_texto = str(recomendacao_raw)
        
        # Montar o item no formato da tabela
        item = {
            "elemento": icone_nome,
            "nc": f"NC-{rec.get('nc_sequencial', '000')}",
            "classificacao": classificacao_texto,
            "recomendacao": recomendacao_texto,  # ← MUDOU: agora garante que é string
            "riscos_mitigados": "",
            "prioridade": rec.get('prioridade', 'amarelo')
        }
        itens.append(item)
    
    # Paginar dinamicamente: 4 itens se alguma recomendação > 120 chars, senão 5
    slides_criados = 0
    i = 0
    while i < len(itens):
        candidatos = itens[i:i + 5]
        tem_longa = any(len(item['recomendacao']) > 120 for item in candidatos)
        lote = itens[i:i + 4] if tem_longa else itens[i:i + 5]
        criar_slide_recomendacoes(pres, dados, lote)
        i += len(lote)
        slides_criados += 1

    print(f"✅ {slides_criados} slides de Recomendações criados!", file=sys.stderr, flush=True)