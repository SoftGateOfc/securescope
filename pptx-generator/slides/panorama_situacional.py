#!/usr/bin/env python3
# -*- coding: utf-8 -*-
import io
import sys
import json
import os
import base64
from io import BytesIO
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from utils.processar_dados import fix_encoding


def add_base64_image(slide, base64_str, left, top, width, height):
    """Adiciona imagem base64 ao slide"""
    try:
        if base64_str and base64_str.startswith('data:image'):
            base64_data = base64_str.split(',')[1]
        else:
            base64_data = base64_str
        
        image_data = base64.b64decode(base64_data)
        image_stream = BytesIO(image_data)
        
        slide.shapes.add_picture(
            image_stream,
            Inches(left),
            Inches(top),
            width=Inches(width),
            height=Inches(height)
        )
        return True
    except Exception as e:
        print(f"Erro ao adicionar imagem base64: {e}", file=sys.stderr)
        return False


def add_local_image(slide, image_path, left, top, width, height):
    """Adiciona imagem local ao slide"""
    try:
        if os.path.exists(image_path):
            slide.shapes.add_picture(
                image_path, 
                Inches(left), 
                Inches(top),
                width=Inches(width), 
                height=Inches(height)
            )
            return True
        else:
            print(f"Imagem não encontrada: {image_path}", file=sys.stderr, flush=True)
            return False
    except Exception as e:
        print(f"Erro ao adicionar imagem local: {e}", file=sys.stderr, flush=True)
        return False


def criar_cabecalho_slide(slide, pres, numero, titulo):
    """Cria o cabeçalho padrão com paralelogramo azul e cinza"""
    
    # === PARALELOGRAMO AZUL (número) === #
    para_azul = slide.shapes.add_shape(
        MSO_SHAPE.PARALLELOGRAM,
        Inches(0.2),    # Left
        Inches(0.20),   # Top (era 0.45)
        Inches(0.5),    # Width
        Inches(0.5)     # Height
    )
    
    para_azul.fill.solid()
    para_azul.fill.fore_color.rgb = RGBColor(30, 115, 190)  # Azul (era 0, 102, 204)
    para_azul.line.fill.background()
    
    # Adicionar número
    tf = para_azul.text_frame
    tf.text = str(numero)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # === PARALELOGRAMO CINZA (título) === #
    para_cinza = slide.shapes.add_shape(
        MSO_SHAPE.PARALLELOGRAM,
        Inches(0.85),   # Left (logo após o paralelogramo azul)
        Inches(0.20),   # Top (era 0.45)
        Inches(3.5),    # Width (era 3.5)
        Inches(0.5)     # Height
    )
    
    para_cinza.fill.solid()
    para_cinza.fill.fore_color.rgb = RGBColor(70, 70, 70)  # Cinza (era 64, 64, 64)
    para_cinza.line.fill.background()
    
    # Adicionar título
    tf = para_cinza.text_frame
    tf.text = titulo
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.font.name = "Arial"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

def gerar_panorama_situacional(pres, dados):
    """Gera o slide de Panorama Situacional"""
    print("Criando slide de Panorama Situacional...", file=sys.stderr)
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    
    # === FUNDO BRANCO === #
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # === CABEÇALHO (Paralelogramo azul + Título) === #
    criar_cabecalho_slide(slide, pres, 3, "PANORAMA SITUACIONAL")
    
    # === DADOS DO RELATÓRIO === #
    dados_info = dados.get("dados", {})
    imagens_info = dados.get("imagens", {})
    
    empresa = fix_encoding(dados.get("dados", {}).get("nome_empresa", ""))
    nome_local = fix_encoding(dados.get("dados", {}).get("localizacao_analise", ""))
    panorama_texto = fix_encoding(dados.get("dados",{}).get("panorama",""))
   
    imagem_area = imagens_info.get("imagem_area", "")
    
    # === TÍTULO "NOME DO LOCAL" (centralizado, acima da imagem) === #
    titulo_local_box = slide.shapes.add_textbox(
        Inches(5.2),   # Alinhado com a borda esquerda da imagem
        Inches(1.1),  # Subido
        Inches(4.5),   # Mesma largura da imagem
        Inches(0.4)  # Altura reduzida
    )
    tf = titulo_local_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = empresa.upper()
    p.font.name = "Arial"
    p.font.size = Pt(18)  # Tamanho reduzido de 24 para 18
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER
    
    # === IMAGEM DO LOCAL (metade direita do slide) === #
    if imagem_area:

         # === SOMBRA DA IMAGEM === #
        sombra = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.2),   
        Inches(1.55),
        Inches(4.6),
        Inches(3.8)
    )

        sombra.fill.solid()
        sombra.fill.fore_color.rgb = RGBColor(255, 255, 255)  # Sem preenchimento
        sombra.line.fill.solid()  # Ativa a borda

        # Cor da borda (escura para dar a impressão de sombra)
        sombra.line.width = Pt(1)  # Largura da borda
        sombra.line.blur = Pt(5)  # Efeito de suavização da borda (ou difusão)
            # Adicionar imagem centralizada no lado direito
        add_base64_image(
            slide,
            imagem_area,
            left=5.2,   # Metade direita
            top=1.6,    # Subido
            width=4.5,
            height=3.7
        )
    else:
        # Placeholder se não houver imagem
        placeholder = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(5.2),
            Inches(1.4),  # Subido
            Inches(4.5),
            Inches(3.2)
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(200, 200, 200)
        placeholder.line.color.rgb = RGBColor(0, 102, 204)
        placeholder.line.width = Pt(2)
        
        tf = placeholder.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "IMAGEM DO LOCAL"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # === LADO ESQUERDO - DADOS DO PANORAMA === #
    
    # 1. TEXTO PANORAMA SITUACIONAL - EXPOSIÇÃO AO RISCO
    if panorama_texto:
        texto_box = slide.shapes.add_textbox(
            Inches(1.0),  # Movido 0.4" para direita (era 0.6)
            Inches(1.1),  # Subido de 1.8 para 1.4
            Inches(3.5),  # Reduzido (era 3.9)
            Inches(0.7)  # Altura reduzida (era 1.0)
        )
        tf = texto_box.text_frame
        tf.clear()
        tf.word_wrap = True
        
        # Título da seção
        p_titulo = tf.paragraphs[0]
        p_titulo.text = "EXPOSIÇÃO AO RISCO"
        p_titulo.font.name = "Arial"
        p_titulo.font.size = Pt(12)
        p_titulo.font.bold = True
        p_titulo.font.color.rgb = RGBColor(0, 51, 102)
        p_titulo.alignment = PP_ALIGN.LEFT
        
        # Texto do panorama
        p_texto = tf.add_paragraph()
        p_texto.text = panorama_texto
        p_texto.font.name = "Arial"
        p_texto.font.size = Pt(10)
        p_texto.font.color.rgb = RGBColor(80, 80, 80)
        p_texto.alignment = PP_ALIGN.LEFT
        p_texto.space_before = Pt(6)
    
    # 2. OBSERVAÇÕES DA LOCALIDADE / REFERÊNCIAS PRÓXIMAS
    ref_texto_box = slide.shapes.add_textbox(
        Inches(1.0),  # Movido 0.4" para direita (era 0.6)
        Inches(2.3),  # Aproximado do texto anterior (era 2.5)
        Inches(3.5),  # Ajustado (era 4.0)
        Inches(0.8)  # Altura reduzida (era 1.2)
    )
    tf = ref_texto_box.text_frame
    tf.clear()
    tf.word_wrap = True
    
    # ÍCONE DE CRIMINALIDADE (ao lado esquerdo do título)
    current_dir = os.path.dirname(__file__)  # pptx-generator/slides/
    parent_dir = os.path.dirname(current_dir)  # pptx-generator/
    icon_criminalidade_path = os.path.join(parent_dir, "images", "ICON_CRIMINALIDADE.png")
    
    add_local_image(
        slide,
        icon_criminalidade_path,
        left=0.3,
        top=2.3,  # Ajustado para acompanhar o texto (era 2.5)
        width=0.6,
        height=0.6
    )
    
    # Título da seção
    p_titulo = tf.paragraphs[0]
    p_titulo.text = "OBSERVAÇÕES DA LOCALIDADE (DADOS, TAXAS DE CRIMINALIDADE ETC.)"
    p_titulo.font.name = "Arial"
    p_titulo.font.size = Pt(9)
    p_titulo.font.bold = True
    p_titulo.font.color.rgb = RGBColor(0, 51, 102)
    p_titulo.alignment = PP_ALIGN.LEFT
    
    # Se tiver lista de referências
    
    p = tf.add_paragraph()
    p.text = "Altere esses dados"
    p.font.name = "Arial"
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(80, 80, 80)
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(3)
   
    # 3. LOCALIZAÇÃO (NOME DO LOCAL + ENDEREÇO)
    loc_texto_box = slide.shapes.add_textbox(
        Inches(1.0),  # Movido 0.4" para direita (era 0.6)
        Inches(3.3),  # Aproximado do texto anterior (era 3.8)
        Inches(3.5),  # Ajustado (era 4.0)
        Inches(0.8)  # Altura reduzida (era 1.0)
    )
    tf = loc_texto_box.text_frame
    tf.clear()
    tf.word_wrap = True
    
    # ÍCONE DE LOCALIZAÇÃO (ao lado esquerdo do título)
    icon_localizacao_path = os.path.join(parent_dir, "images", "ICON_LOCALIZAÇÃO.png")
    
    add_local_image(
        slide,
        icon_localizacao_path,
        left=0.3,
        top=3.3,  # Ajustado para acompanhar o texto (era 3.8)
        width=0.6,
        height=0.6
    )
    
    # Título LOCALIZAÇÃO
    p_loc_titulo = tf.paragraphs[0]
    p_loc_titulo.text = "LOCALIZAÇÃO"
    p_loc_titulo.font.name = "Arial"
    p_loc_titulo.font.size = Pt(9)
    p_loc_titulo.font.bold = True
    p_loc_titulo.font.color.rgb = RGBColor(0, 51, 102)
    p_loc_titulo.alignment = PP_ALIGN.LEFT
    
    # Nome do local
    p_nome = tf.add_paragraph()
    p_nome.text = nome_local
    p_nome.font.name = "Arial"
    p_nome.font.size = Pt(9)
    p_nome.font.bold = True
    p_nome.font.color.rgb = RGBColor(80, 80, 80)
    p_nome.alignment = PP_ALIGN.LEFT
    p_nome.space_before = Pt(6)
    
    # Endereço (placeholder)
    p_endereco = tf.add_paragraph()
    p_endereco.text = " Adicione o Endereço"
    p_endereco.font.name = "Arial"
    p_endereco.font.size = Pt(9)
    p_endereco.font.color.rgb = RGBColor(80, 80, 80)
    p_endereco.alignment = PP_ALIGN.LEFT
    p_endereco.space_before = Pt(3)
    
    caixa_azul = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.3),    # Começa do lado esquerdo
    Inches(4.3),    # Abaixo da seção de localização
    Inches(4.5),    # Até um pouco antes da metade da página
    Inches(1.2)     # Altura da caixa
)

# Fundo branco
    caixa_azul.fill.solid()
    caixa_azul.fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Borda azul (0,102,204)
    caixa_azul.line.fill.solid()
    caixa_azul.line.color.rgb = RGBColor(0, 102, 204)
    caixa_azul.line.width = Pt(2) 

    icone_width = 0.45
    icone_height = 0.45

    texto_offset_x = 0.45
    texto_largura = 1.2
    texto_altura = 0.3

    # Caminho da pasta de imagens
    icones_path = os.path.join(parent_dir, "images")

    # Ícones e textos
    icones_superior = [
        ("policia_militar.png",  "Polícia Militar"),
        ("policia_civil.png",    "Polícia Civil"),
        ("policia_federal.png",  "Polícia Federal")
    ]

    icones_inferior = [
        ("bombeiro.png",  "Corpo de Bombeiros"),
        ("ICON_PORTO.png", "Porto")
    ]

    base_top_superior = 4.38
    base_top_inferior = 4.95

    posicoes_superior = [0.45, 1.8, 3.15]
    posicoes_inferior = [1.0, 2.8]

    # === Linha Superior (3 ícones) === #
    for i, (arquivo, titulo) in enumerate(icones_superior):
        caminho = os.path.join(icones_path, arquivo)

        add_local_image(
            slide,
            caminho,
            left=posicoes_superior[i],
            top=base_top_superior,
            width=icone_width,
            height=icone_height
        )

        # Texto ao lado
        tb = slide.shapes.add_textbox(
            Inches(posicoes_superior[i] + texto_offset_x),
            Inches(base_top_superior + 0.05),
            Inches(texto_largura),
            Inches(texto_altura)
        )
        tf = tb.text_frame
        tf.clear()

        p1 = tf.paragraphs[0]
        p1.text = titulo
        p1.font.name = "Arial"
        p1.font.size = Pt(8)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(0, 51, 102)

        p2 = tf.add_paragraph()
        p2.text = "_ Km"
        p2.font.name = "Arial"
        p2.font.size = Pt(8)
        p2.font.color.rgb = RGBColor(80, 80, 80)

    # === Linha Inferior (2 ícones) === #
    for i, (arquivo, titulo) in enumerate(icones_inferior):
        caminho = os.path.join(icones_path, arquivo)

        add_local_image(
            slide,
            caminho,
            left=posicoes_inferior[i],
            top=base_top_inferior,
            width=icone_width,
            height=icone_height
        )

        tb = slide.shapes.add_textbox(
            Inches(posicoes_inferior[i] + texto_offset_x),
            Inches(base_top_inferior + 0.05),
            Inches(texto_largura),
            Inches(texto_altura)
        )
        tf = tb.text_frame
        tf.clear()

        p1 = tf.paragraphs[0]
        p1.text = titulo
        p1.font.name = "Arial"
        p1.font.size = Pt(8)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(0, 51, 102)

        p2 = tf.add_paragraph()
        p2.text = "_ Km"
        p2.font.name = "Arial"
        p2.font.size = Pt(8)
        p2.font.color.rgb = RGBColor(80, 80, 80)
    
    print("Slide de Panorama Situacional criado!", file=sys.stderr)
    return slide