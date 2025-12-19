#!/usr/bin/env python3
# -*- coding: utf-8 -*-
import sys
import os
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import base64
from io import BytesIO

# Importar funções de processamento
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..'))
from utils.processar_dados import nomear_nivel_vulnerabilidade


def add_base64_image(slide, base64_str, left, top, width, height):
    """Adiciona imagem base64 ao slide"""
    try:
        if ',' in base64_str:
            base64_str = base64_str.split(',')[1]
        
        img_data = base64.b64decode(base64_str)
        stream = BytesIO(img_data)
        slide.shapes.add_picture(
            stream, Inches(left), Inches(top),
            width=Inches(width), height=Inches(height)
        )
        return True
    except Exception as e:
        print(f"❌ Erro ao adicionar imagem: {e}", file=sys.stderr, flush=True)
        return False


def mapear_criticidade_texto(criticidade):
    """Mapeia criticidade para texto descritivo"""
    mapa = {
        'vermelho-escuro': 'Muito Alto',
        'laranja': 'Alto',
        'amarelo': 'Médio',
        'verde-escuro': 'Baixo',
        'verde-claro': 'Muito Baixo'
    }
    return mapa.get(criticidade, 'Não Classificado')


def mapear_prioridade_texto(prioridade):
    """Mapeia prioridade para texto descritivo"""
    mapa = {
        'vermelho-escuro': 'Longo Prazo',
        'amarelo': 'Médio Prazo',
        'verde-claro': 'Curto Prazo'
    }
    return mapa.get(prioridade, 'Não Classificado')


def extrair_pilar_nome(pilar_path):
    """Extrai nome do pilar do path da imagem"""
    pilar_lower = pilar_path.lower()
    
    if 'pessoas' in pilar_lower:
        return 'PESSOAS'
    elif 'tecnologia' in pilar_lower:
        return 'TECNOLOGIA'
    elif 'processos' in pilar_lower:
        return 'PROCESSOS'
    elif 'informacao' in pilar_lower or 'informação' in pilar_lower:
        return 'INFORMAÇÃO'
    elif 'gestao' in pilar_lower or 'gestão' in pilar_lower:
        return 'GESTÃO'
    
    return 'PILAR'


def criar_caixa_informacoes(slide, vulnerabilidade, left, top, width, height):
    """
    Cria caixa de informações com dados da vulnerabilidade
    """
    # === CAIXA PRINCIPAL === #
    caixa = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height)
    )
    
    # Fundo branco
    caixa.fill.solid()
    caixa.fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Borda azul
    caixa.line.fill.solid()
    caixa.line.color.rgb = RGBColor(30, 115, 190)
    caixa.line.width = Pt(2)
    
    # === EXTRAIR DADOS === #
    pilar_path = vulnerabilidade.get('pilar', '')
    pilar_nome = extrair_pilar_nome(pilar_path)
    topicos = vulnerabilidade.get('topicos', 'Sem descrição').strip()
    titulo_pergunta = vulnerabilidade.get('titulo_pergunta', 'Sem título').strip()
    vulnerabilidade_nivel = vulnerabilidade.get('vulnerabilidade', 0)
    adequacao_texto = nomear_nivel_vulnerabilidade(vulnerabilidade_nivel)
    criticidade = vulnerabilidade.get('criticidade', '')
    risco_texto = mapear_criticidade_texto(criticidade)
    prioridade = vulnerabilidade.get('prioridade', '')
    prioridade_texto = mapear_prioridade_texto(prioridade)
    
    # === 🔥 ÍCONE DO PILAR NO TOPO === #
    # Mapear caminho do ícone
    current_dir = os.path.dirname(__file__)
    parent_dir = os.path.dirname(current_dir)
    images_dir = os.path.join(parent_dir, "images")
    
    # Pegar o nome correto do ícone baseado no pilar
    icone_map = {
        'PESSOAS': 'pessoas.png',
        'TECNOLOGIA': 'tecnologia.png',
        'PROCESSOS': 'processos.png',
        'INFORMAÇÃO': 'informacoes.png',
        'GESTÃO': 'gestao.png'
    }
    
    icone_nome = icone_map.get(pilar_nome, 'pessoas.png')
    icone_path = os.path.join(images_dir, icone_nome)
    
    # Adicionar ícone do pilar (centralizado no topo da caixa)
    icone_size = 0.4
    icone_left = left + (width - icone_size) / 2  # Centralizar
    icone_top = top + 0.1
    
    try:
        if os.path.exists(icone_path):
            slide.shapes.add_picture(
                icone_path,
                Inches(icone_left),
                Inches(icone_top),
                width=Inches(icone_size),
                height=Inches(icone_size)
            )
    except Exception as e:
        print(f"⚠️ Erro ao adicionar ícone do pilar: {e}", file=sys.stderr)
    
    # Nome do pilar abaixo do ícone
    pilar_text_box = slide.shapes.add_textbox(
        Inches(left + 0.1),
        Inches(icone_top + icone_size + 0.05),
        Inches(width - 0.2),
        Inches(0.25)
    )
    tf_pilar = pilar_text_box.text_frame
    tf_pilar.clear()
    p_pilar = tf_pilar.paragraphs[0]
    p_pilar.text = pilar_nome
    p_pilar.font.name = "Arial"
    p_pilar.font.size = Pt(12)
    p_pilar.font.bold = True
    p_pilar.font.color.rgb = RGBColor(0, 51, 102)
    p_pilar.alignment = PP_ALIGN.CENTER
    
    # === TEXTBOX COM INFORMAÇÕES (abaixo do pilar) === #
    texto_top = icone_top + icone_size + 0.35  # Abaixo do nome do pilar
    
    text_box = slide.shapes.add_textbox(
        Inches(left + 0.15),
        Inches(texto_top),
        Inches(width - 0.3),
        Inches(height - (texto_top - top) - 0.15)
    )
    
    tf = text_box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    
    # Cores
    cor_label = RGBColor(0, 51, 102)
    cor_valor = RGBColor(30, 115, 190)
    
    # === 1. PERGUNTA: (inline) === #
    p_pergunta = tf.paragraphs[0]
    
    # Run do label "Pergunta:"
    run_label = p_pergunta.add_run()
    run_label.text = "Pergunta: "
    run_label.font.name = "Arial"
    run_label.font.size = Pt(10)
    run_label.font.bold = True
    run_label.font.color.rgb = cor_label
    
    # Run do valor (inline)
    run_valor = p_pergunta.add_run()
    run_valor.text = titulo_pergunta
    run_valor.font.name = "Arial"
    run_valor.font.size = Pt(9)
    run_valor.font.bold = False
    run_valor.font.color.rgb = cor_valor
    
    p_pergunta.space_after = Pt(10)
    
    # === 2. TÓPICO: (inline) === #
    p_topico = tf.add_paragraph()
    
    run_label = p_topico.add_run()
    run_label.text = "Tópico: "
    run_label.font.name = "Arial"
    run_label.font.size = Pt(10)
    run_label.font.bold = True
    run_label.font.color.rgb = cor_label
    
    run_valor = p_topico.add_run()
    run_valor.text = topicos
    run_valor.font.name = "Arial"
    run_valor.font.size = Pt(9)
    run_valor.font.bold = False
    run_valor.font.color.rgb = cor_valor
    
    p_topico.space_after = Pt(10)
    
    # === 3. ADEQUAÇÃO: (inline) === #
    p_adequacao = tf.add_paragraph()
    
    run_label = p_adequacao.add_run()
    run_label.text = "Adequação: "
    run_label.font.name = "Arial"
    run_label.font.size = Pt(10)
    run_label.font.bold = True
    run_label.font.color.rgb = cor_label
    
    run_valor = p_adequacao.add_run()
    run_valor.text = adequacao_texto
    run_valor.font.name = "Arial"
    run_valor.font.size = Pt(9)
    run_valor.font.bold = False
    run_valor.font.color.rgb = cor_valor
    
    p_adequacao.space_after = Pt(10)
    
    # === 4. RISCO: (inline) === #
    p_risco = tf.add_paragraph()
    
    run_label = p_risco.add_run()
    run_label.text = "Risco: "
    run_label.font.name = "Arial"
    run_label.font.size = Pt(10)
    run_label.font.bold = True
    run_label.font.color.rgb = cor_label
    
    run_valor = p_risco.add_run()
    run_valor.text = risco_texto
    run_valor.font.name = "Arial"
    run_valor.font.size = Pt(9)
    run_valor.font.bold = False
    run_valor.font.color.rgb = cor_valor
    
    p_risco.space_after = Pt(10)
    
    # === 5. PRIORIDADE: (inline) === #
    p_prioridade = tf.add_paragraph()
    
    run_label = p_prioridade.add_run()
    run_label.text = "Prioridade: "
    run_label.font.name = "Arial"
    run_label.font.size = Pt(10)
    run_label.font.bold = True
    run_label.font.color.rgb = cor_label
    
    run_valor = p_prioridade.add_run()
    run_valor.text = prioridade_texto
    run_valor.font.name = "Arial"
    run_valor.font.size = Pt(9)
    run_valor.font.bold = False
    run_valor.font.color.rgb = cor_valor
    
    return caixa


def criar_slide_vulnerabilidade_foto(pres, dados, vulnerabilidade):
    """Cria um slide de vulnerabilidade com foto"""
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    
    # === FUNDO BRANCO === #
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # === CABEÇALHO === #
    para_azul = slide.shapes.add_shape(
        MSO_SHAPE.PARALLELOGRAM,
        Inches(0.2), Inches(0.2), Inches(0.5), Inches(0.5)
    )
    para_azul.fill.solid()
    para_azul.fill.fore_color.rgb = RGBColor(30, 115, 190)
    para_azul.line.fill.background()
    
    tf_azul = para_azul.text_frame
    tf_azul.text = "8"
    p = tf_azul.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    tf_azul.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    para_cinza = slide.shapes.add_shape(
        MSO_SHAPE.PARALLELOGRAM,
        Inches(0.85), Inches(0.2), Inches(5.8), Inches(0.5)
    )
    para_cinza.fill.solid()
    para_cinza.fill.fore_color.rgb = RGBColor(70, 70, 70)
    para_cinza.line.fill.background()
    
    tf_cinza = para_cinza.text_frame
    tf_cinza.text = "VULNERABILIDADES ESPECÍFICAS"
    tf_cinza.margin_left = Inches(0.2)
    tf_cinza.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_cinza.paragraphs[0].font.bold = True
    tf_cinza.paragraphs[0].font.size = Pt(20)
    tf_cinza.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    subtitulo = slide.shapes.add_textbox(
        Inches(2.5), Inches(0.75), Inches(5.0), Inches(0.55)
    )
    p = subtitulo.text_frame.paragraphs[0]
    p.text = "Áreas Restritas"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 38, 77)
    p.alignment = PP_ALIGN.CENTER

    # === CAIXA DE INFORMAÇÕES === #
    criar_caixa_informacoes(slide, vulnerabilidade, 0.3, 1.7, 3.0, 2.9)
    
    # === FOTO === #
    foto_base64 = vulnerabilidade.get('foto_base64', None)
    
    if foto_base64:
        print(f"  📸 Adicionando foto NC-{vulnerabilidade.get('nc', '???')}", file=sys.stderr, flush=True)
        sucesso = add_base64_image(slide, foto_base64, 3.5, 1.4, 6.0, 3.2)
        
        if not sucesso:
            placeholder = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(3.5), Inches(1.0), Inches(6.0), Inches(3.5)
            )
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = RGBColor(240, 240, 240)
            placeholder.line.color.rgb = RGBColor(200, 0, 0)
            placeholder.line.width = Pt(2)
    else:
        texto_sem_foto = slide.shapes.add_textbox(
            Inches(3.5), Inches(2.5), Inches(6.0), Inches(1.0)
        )
        tf_sem = texto_sem_foto.text_frame
        p_sem = tf_sem.paragraphs[0]
        p_sem.text = "Sem registro fotográfico"
        p_sem.font.name = "Arial"
        p_sem.font.size = Pt(14)
        p_sem.font.italic = True
        p_sem.font.color.rgb = RGBColor(150, 150, 150)
        p_sem.alignment = PP_ALIGN.CENTER
        tf_sem.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # === LOGO === #
    logo = dados.get("imagens", {}).get("logo_empresa")
    if logo:
        try:
            if ',' in logo:
                logo = logo.split(',')[1]
            stream = BytesIO(base64.b64decode(logo))
            slide.shapes.add_picture(stream, Inches(9.0), Inches(4.7), width=Inches(0.8))
        except Exception as e:
            print(f"❌ Erro logo: {e}", file=sys.stderr, flush=True)
    
    return slide


def gerar_vulnerabilidades_foto(pres, dados):
    """Gera slides de vulnerabilidades com fotos"""
    print("🔄 Processando vulnerabilidades com fotos...", file=sys.stderr, flush=True)
    
    respostas = dados.get('dados_modelo', {}).get('respostas', [])
    print(f"📊 Total de respostas: {len(respostas)}", file=sys.stderr, flush=True)
    
    if not respostas:
        print("⚠️ Nenhuma resposta encontrada", file=sys.stderr, flush=True)
        return
    
    # 🔥 DEBUG: Ver quantas têm foto
    total_vulnerabilidades = 0
    total_com_foto = 0
    
    for r in respostas:
        vuln = int(r.get('vulnerabilidade', 0))
        foto = r.get('foto_base64')
        
        if vuln > 0:
            total_vulnerabilidades += 1
            if foto:
                total_com_foto += 1
                print(f"  ✓ NC-{r.get('nc', '???')} tem foto (tamanho: {len(foto)} chars)", file=sys.stderr, flush=True)
            else:
                print(f"  ✗ NC-{r.get('nc', '???')} SEM foto", file=sys.stderr, flush=True)
    
    print(f"📊 Vulnerabilidades (nivel > 1): {total_vulnerabilidades}", file=sys.stderr, flush=True)
    print(f"📸 Vulnerabilidades COM foto: {total_com_foto}", file=sys.stderr, flush=True)
    
    # Filtrar vulnerabilidades COM foto
    vulnerabilidades_com_foto = [
        r for r in respostas 
        if int(r.get('vulnerabilidade', 0)) > 0 and r.get('foto_base64')
    ]
    
    if len(vulnerabilidades_com_foto) == 0:
        print("⚠️ Nenhuma vulnerabilidade com foto após filtro", file=sys.stderr, flush=True)
        return
    
    print(f"🎯 Criando {len(vulnerabilidades_com_foto)} slides...", file=sys.stderr, flush=True)
    
    # 🔥 Criar um slide para CADA vulnerabilidade
    for idx, vuln in enumerate(vulnerabilidades_com_foto, start=1):
        print(f"  📄 Slide {idx}/{len(vulnerabilidades_com_foto)} - NC-{vuln.get('nc', '???')}", file=sys.stderr, flush=True)
        criar_slide_vulnerabilidade_foto(pres, dados, vuln)
    
    print(f"✅ {len(vulnerabilidades_com_foto)} slides criados!", file=sys.stderr, flush=True)