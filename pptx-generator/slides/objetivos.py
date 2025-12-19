import sys
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement  
from utils.processar_dados import fix_encoding


def add_triangle_top_right(slide, pres):
    """Triângulo branco na quina SUPERIOR DIREITA"""
    width = pres.slide_width
    height = pres.slide_height
    triangle_base = Inches(1.8)
    triangle_height = height / 2.3

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 1, 1)
    sp = shape._element
    spPr = sp.spPr

    for child in list(spPr):
        if "prstGeom" in child.tag:
            spPr.remove(child)

    custGeom = OxmlElement("a:custGeom")
    for tag in ("a:avLst", "a:gdLst", "a:ahLst", "a:cxnLst"):
        custGeom.append(OxmlElement(tag))

    pathLst = OxmlElement("a:pathLst")
    path = OxmlElement("a:path")
    path.set("w", str(width))
    path.set("h", str(height))

    move = OxmlElement("a:moveTo")
    pt = OxmlElement("a:pt")
    pt.set("x", str(width))
    pt.set("y", "0")
    move.append(pt)
    path.append(move)

    ln1 = OxmlElement("a:lnTo")
    p1 = OxmlElement("a:pt")
    p1.set("x", str(width - int(triangle_base)))
    p1.set("y", "0")
    ln1.append(p1)
    path.append(ln1)

    ln2 = OxmlElement("a:lnTo")
    p2 = OxmlElement("a:pt")
    p2.set("x", str(width))
    p2.set("y", str(int(triangle_height)))
    ln2.append(p2)
    path.append(ln2)

    path.append(OxmlElement("a:close"))
    pathLst.append(path)
    custGeom.append(pathLst)
    spPr.append(custGeom)

    shape.left = 0
    shape.top = 0
    shape.width = width
    shape.height = height

    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.fill.background()


def add_triangle_bottom_right(slide, pres):
    """Triângulo branco na quina INFERIOR DIREITA"""
    width = pres.slide_width
    height = pres.slide_height
    triangle_base = Inches(1.8)
    triangle_height = height / 2.3

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 1, 1)
    sp = shape._element
    spPr = sp.spPr

    for child in list(spPr):
        if "prstGeom" in child.tag:
            spPr.remove(child)

    custGeom = OxmlElement("a:custGeom")
    for tag in ("a:avLst", "a:gdLst", "a:ahLst", "a:cxnLst"):
        custGeom.append(OxmlElement(tag))

    pathLst = OxmlElement("a:pathLst")
    path = OxmlElement("a:path")
    path.set("w", str(width))
    path.set("h", str(height))

    move = OxmlElement("a:moveTo")
    pt = OxmlElement("a:pt")
    pt.set("x", str(width))
    pt.set("y", str(height))
    move.append(pt)
    path.append(move)

    ln1 = OxmlElement("a:lnTo")
    p1 = OxmlElement("a:pt")
    p1.set("x", str(width - int(triangle_base)))
    p1.set("y", str(height))
    ln1.append(p1)
    path.append(ln1)

    ln2 = OxmlElement("a:lnTo")
    p2 = OxmlElement("a:pt")
    p2.set("x", str(width))
    p2.set("y", str(height - int(triangle_height)))
    ln2.append(p2)
    path.append(ln2)

    path.append(OxmlElement("a:close"))
    pathLst.append(path)
    custGeom.append(pathLst)
    spPr.append(custGeom)

    shape.left = 0
    shape.top = 0
    shape.width = width
    shape.height = height

    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.fill.background()


def add_triangle_center_isosceles(slide, pres):
    """Triângulo isósceles branco no CENTRO"""
    width = pres.slide_width
    height = pres.slide_height
    
    meio_x = width / 2
    triangle_height = Inches(1)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 1, 1)
    sp = shape._element
    spPr = sp.spPr

    for child in list(spPr):
        if "prstGeom" in child.tag:
            spPr.remove(child)

    custGeom = OxmlElement("a:custGeom")
    for tag in ("a:avLst", "a:gdLst", "a:ahLst", "a:cxnLst"):
        custGeom.append(OxmlElement(tag))

    pathLst = OxmlElement("a:pathLst")
    path = OxmlElement("a:path")
    path.set("w", str(width))
    path.set("h", str(height))

    move = OxmlElement("a:moveTo")
    pt = OxmlElement("a:pt")
    pt.set("x", str(int(meio_x)))
    pt.set("y", "0")
    move.append(pt)
    path.append(move)

    ln1 = OxmlElement("a:lnTo")
    p1 = OxmlElement("a:pt")
    p1.set("x", str(int(meio_x + triangle_height)))
    p1.set("y", str(int(height / 2)))
    ln1.append(p1)
    path.append(ln1)

    ln2 = OxmlElement("a:lnTo")
    p2 = OxmlElement("a:pt")
    p2.set("x", str(int(meio_x)))
    p2.set("y", str(int(height)))
    ln2.append(p2)
    path.append(ln2)

    path.append(OxmlElement("a:close"))
    pathLst.append(path)
    custGeom.append(pathLst)
    spPr.append(custGeom)

    shape.left = 0
    shape.top = 0
    shape.width = width
    shape.height = height

    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(220, 220, 220)
    shape.line.fill.background()


def gerar_objetivos(pres, dados):
    print("Criando slide de objetivos...", file=sys.stderr, flush=True)
    
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    largura_total = pres.slide_width

    largura_faixa = largura_total * 0.50
    posicao_faixa = largura_total - largura_faixa
    
    faixa_azul = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        int(posicao_faixa),
        0,
        int(largura_faixa),
        pres.slide_height
    )
    
    # === GRADIENTE AZUL (escuro embaixo, claro em cima) === #
    fill_faixa = faixa_azul.fill
    fill_faixa.gradient()
    fill_faixa.gradient_angle = 90  # Vertical
    
    # Stop 0: Azul escuro (embaixo)
    stop_dark = fill_faixa.gradient_stops[0]
    stop_dark.position = 0.0
    stop_dark.color.rgb = RGBColor(10, 50, 100)
    
    # Stop 1: Azul claro (em cima)
    stop_light = fill_faixa.gradient_stops[1]
    stop_light.position = 1.0
    stop_light.color.rgb = RGBColor(30, 115, 190)
    
    # Remove borda
    faixa_azul.line.fill.background()
    faixa_azul.line.width = Pt(0)
    
    sp_faixa = faixa_azul._element
    spPr_faixa = sp_faixa.spPr
    for ln_elem in spPr_faixa.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}ln'):
        spPr_faixa.remove(ln_elem)

    add_triangle_top_right(slide, pres)
    add_triangle_bottom_right(slide, pres)
    add_triangle_center_isosceles(slide, pres)

    # === TEXTO 1: "Elevando o nível de segurança" (0.15 inches mais à esquerda) === #
    texto1 = slide.shapes.add_textbox(
        int(posicao_faixa) + Inches(1),  # 0.15 mais à esquerda (era 1.0)
        Inches(0.45),
        Inches(3.15),  # Largura ajustada
        Inches(0.6)
    )
    tf1 = texto1.text_frame
    tf1.word_wrap = True
    tf1.vertical_anchor = MSO_ANCHOR.TOP

    p1 = tf1.paragraphs[0]
    p1.text = "Elevando o nível de segurança"
    p1.font.name = 'Aptos'
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(255, 255, 255)
    p1.alignment = PP_ALIGN.LEFT

    # === TEXTO 2 e 3: Parágrafos na posição original === #

    nome_empresa = fix_encoding(dados.get("dados", {}).get("nome_empresa", ""))
    localizacao_analise  = fix_encoding(dados.get("dados", {}).get("localizacao_analise", ""))

    texto2 = slide.shapes.add_textbox(
        int(posicao_faixa) + Inches(1.0),  # Posição original
        Inches(1.15),  # Abaixo do texto1
        Inches(3.0),
        Inches(4.0)
    )
    tf2 = texto2.text_frame
    tf2.word_wrap = True
    tf2.vertical_anchor = MSO_ANCHOR.TOP

    p2 = tf2.paragraphs[0]
    p2.text = f"Realizar análise de riscos em {localizacao_analise}, propondo soluções priorizadas para mitigar vulnerabilidades identificadas."
    p2.font.name = 'Aptos'
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(255, 255, 255)
    p2.alignment = PP_ALIGN.LEFT
    p2.space_after = Pt(100)

    p3 = tf2.add_paragraph()
    p3.text = f"Os resultados permitirão a {nome_empresa} planejar ações corretivas e prevenir incidentes que afetem pessoas, ativos e a operação de forma contínua e eficaz."
    p3.font.name = 'Aptos'
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(255, 255, 255)
    p3.alignment = PP_ALIGN.LEFT

    titulo_box = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(0.5),
        Inches(4),
        Inches(0.6)
    )
    tf_titulo = titulo_box.text_frame
    tf_titulo.text = "OBJETIVOS"

    p_titulo = tf_titulo.paragraphs[0]
    p_titulo.font.name = 'Aptos'
    p_titulo.font.size = Pt(32)
    p_titulo.font.bold = True
    p_titulo.font.color.rgb = RGBColor(0, 112, 192)

    linha = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5),
        Inches(1.13),
        Inches(3.6),
        Inches(0.015)
    )
    linha.fill.solid()
    linha.fill.fore_color.rgb = RGBColor(0, 0, 0)
    linha.line.fill.background()

    topicos = [
        "Identificação de Risco e Vulnerabilidade",
        "Definição de Recomendações",
        "Classificação por Criticidade e Prioridade",
        "Consolidação em Relatório Executivo"
    ]

    start_top = Inches(2.2)
    spacing = Inches(0.75)
    circle_left = Inches(0.45)
    text_left = Inches(1.15)

    box_width = Inches(3.4)
    box_height = Inches(0.45)
    blue_color = RGBColor(0, 112, 192)

    for i, topico_texto in enumerate(topicos):
        topic_top = start_top + (i * spacing)

        circulo = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            circle_left,
            topic_top,
            Inches(0.38),
            Inches(0.38)
        )
        circulo.fill.solid()
        circulo.fill.fore_color.rgb = blue_color
        circulo.line.fill.background()

        tf_c = circulo.text_frame
        tf_c.text = str(i + 1)
        tf_c.vertical_anchor = MSO_ANCHOR.MIDDLE
        pc = tf_c.paragraphs[0]
        pc.alignment = PP_ALIGN.CENTER
        pc.font.name = 'Aptos'
        pc.font.size = Pt(16)
        pc.font.bold = True
        pc.font.color.rgb = RGBColor(255, 255, 255)

        caixa = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            text_left,
            topic_top - Inches(0.05),
            box_width,
            box_height
        )
        caixa.fill.solid()
        caixa.fill.fore_color.rgb = RGBColor(255, 255, 255)
        caixa.line.color.rgb = blue_color
        caixa.line.width = Pt(1.2)

        tfb = caixa.text_frame
        tfb.text = topico_texto
        tfb.margin_left = Inches(0.1)
        tfb.margin_right = Inches(0.1)
        tfb.vertical_anchor = MSO_ANCHOR.MIDDLE

        pb = tfb.paragraphs[0]
        pb.font.name = 'Aptos'
        pb.font.size = Pt(12.5)
        pb.font.bold = True
        pb.font.color.rgb = blue_color

    logo = dados.get("imagens", {}).get("logo_empresa")
    if logo:
        try:
            from io import BytesIO
            import base64

            if ',' in logo:
                logo = logo.split(',')[1]

            img_data = base64.b64decode(logo)
            stream = BytesIO(img_data)

            slide.shapes.add_picture(
                stream,
                int(posicao_faixa) + Inches(3.8),
                Inches(4.5),
                width=Inches(0.9)
            )
        except Exception as e:
            print(f"Erro ao adicionar logo: {e}", file=sys.stderr, flush=True)

    print("Slide de objetivos criado!", file=sys.stderr, flush=True)
    return slide