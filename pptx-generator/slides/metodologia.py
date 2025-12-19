#!/usr/bin/env python3
# -*- coding: utf-8 -*-
import sys
import os
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


def add_local_image(slide, image_path, left, top, width, height):
    """Adiciona imagem local ao slide"""
    try:
        if os.path.exists(image_path):
            slide.shapes.add_picture(
                image_path, Inches(left), Inches(top),
                width=Inches(width), height=Inches(height)
            )
            return True
        else:
            print(f"Imagem não encontrada: {image_path}", file=sys.stderr)
            return False
    except Exception as e:
        print(f"Erro ao adicionar imagem local: {e}", file=sys.stderr)
        return False


def add_texto_pilar(slide, titulo, descricao, left, top, width, align=PP_ALIGN.LEFT, cor_titulo=(0,51,102)):
    """Adiciona texto de título e descrição para um pilar"""
    try:
        text_box = slide.shapes.add_textbox(
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(1.0)
        )
        
        tf = text_box.text_frame
        tf.clear()
        tf.word_wrap = True
        
        # Título (agora com cor customizável)
        p_titulo = tf.paragraphs[0]
        p_titulo.text = titulo
        p_titulo.font.name = "Arial"
        p_titulo.font.size = Pt(11)
        p_titulo.font.bold = True
        p_titulo.font.color.rgb = RGBColor(*cor_titulo)
        p_titulo.alignment = align
        
        # Descrição (cinza padrão)
        p_desc = tf.add_paragraph()
        p_desc.text = descricao
        p_desc.font.name = "Arial"
        p_desc.font.size = Pt(9)
        p_desc.font.bold = False
        p_desc.font.color.rgb = RGBColor(70,70,70)
        p_desc.alignment = align
        p_desc.space_before = Pt(2)
        
        return True
    except Exception as e:
        print(f"Erro ao adicionar texto do pilar: {e}", file=sys.stderr)
        return False


def gerar_metodologia(pres, dados):
    """Gera o slide de Metodologia (Página 4)"""
    print("Criando slide de Metodologia...", file=sys.stderr)
    slide = pres.slides.add_slide(pres.slide_layouts[6])

    # === FUNDO BRANCO === #
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # === LOGO EMPRESA (canto inferior direito) === #
    logo = dados.get("imagens", {}).get("logo_empresa")
    if logo:
        try:
            from io import BytesIO
            import base64
            
            if ',' in logo:
                logo = logo.split(',')[1]
            
            img_data = base64.b64decode(logo)
            stream = BytesIO(img_data)
            
            # Calcular posição no canto inferior direito
            slide.shapes.add_picture(
                stream,
                Inches(9.0),   # Canto direito (slide tem 10" de largura)
                Inches(4.5),   # Canto inferior (slide tem 5.625" de altura)
                width=Inches(0.9)
            )
        except Exception as e:
            print(f"Erro ao adicionar logo: {e}", file=sys.stderr, flush=True)

    # === PARALELOGRAMO AZUL COM NÚMERO "2" === #
    para_azul = slide.shapes.add_shape(
        MSO_SHAPE.PARALLELOGRAM,
        Inches(0.2),    # Posição à esquerda
        Inches(0.20),   # Alinhado com o título
        Inches(0.5),    # Largura
        Inches(0.5)     # Altura
    )
    para_azul.fill.solid()
    para_azul.fill.fore_color.rgb = RGBColor(30, 115, 190)  # Azul
    para_azul.line.fill.background()
    
    # Adicionar número "2" no paralelogramo azul
    tf_azul = para_azul.text_frame
    tf_azul.text = "2"
    tf_azul.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_azul = tf_azul.paragraphs[0]
    p_azul.alignment = PP_ALIGN.CENTER
    p_azul.font.name = "Arial"
    p_azul.font.size = Pt(20)
    p_azul.font.bold = True
    p_azul.font.color.rgb = RGBColor(255, 255, 255)  # Branco

    # === PARALELOGRAMO CINZA COM TÍTULO "METODOLOGIA" === #
    para_cinza = slide.shapes.add_shape(
        MSO_SHAPE.PARALLELOGRAM,
        Inches(0.85),   # Logo após o paralelogramo azul
        Inches(0.20),   # Alinhado verticalmente
        Inches(3.5),    # Largura para o texto
        Inches(0.5)     # Altura
    )
    para_cinza.fill.solid()
    para_cinza.fill.fore_color.rgb = RGBColor(70, 70, 70)  # Cinza
    para_cinza.line.fill.background()
    
    # Adicionar texto "METODOLOGIA" no paralelogramo cinza
    tf_cinza = para_cinza.text_frame
    tf_cinza.text = "METODOLOGIA"
    tf_cinza.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_cinza.margin_left = Inches(0.2)
    p_cinza = tf_cinza.paragraphs[0]
    p_cinza.alignment = PP_ALIGN.LEFT
    p_cinza.font.name = "Arial"
    p_cinza.font.size = Pt(14)
    p_cinza.font.bold = True
    p_cinza.font.color.rgb = RGBColor(255, 255, 255)  # Branco

    # === SUBTÍTULO "Metodologia Totalizante" (subido 1.5", centralizado) === #
    subtitulo_box = slide.shapes.add_textbox(
        Inches(0.2),
        Inches(2.0),  
        Inches(4.0),
        Inches(0.8)   # Aumentado para caber as duas linhas
    )
    tf = subtitulo_box.text_frame
    tf.clear()
    tf.word_wrap = True
    
    # Primeira linha: "Metodologia" em azul escuro e negrito
    p1 = tf.paragraphs[0]
    p1.text = "Metodologia"
    p1.font.name = "Arial"
    p1.font.size = Pt(22)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(0, 51, 102)  # Azul mais escuro
    p1.alignment = PP_ALIGN.CENTER
    
    # Segunda linha: "Totalizante" mantém azul claro
    p2 = tf.add_paragraph()
    p2.text = "Totalizante"
    p2.font.name = "Arial"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(0, 102, 204)  # Azul claro (como era antes)
    p2.alignment = PP_ALIGN.CENTER

    # === TEXTO DESCRITIVO (subido 1.5", com 0.5" de espaçamento) === #
    texto_box = slide.shapes.add_textbox(
        Inches(0.2),
        Inches(3.3),  # 1.8 + 0.8 (altura subtítulo) + 0.5 (espaçamento) = 3.1
        Inches(4.0),
        Inches(1.0)
    )
    tf = texto_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Soluções avançadas em classificação de riscos para inteligência do seu negócio."
    p.font.name = "Arial"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0, 102, 204)  
    p.alignment = PP_ALIGN.CENTER
    tf.word_wrap = True

    # === IMAGEM DOS 5 PILARES (Arte_metodologia.png) === #
    current_dir = os.path.dirname(__file__)  # pptx-generator/slides/
    parent_dir = os.path.dirname(current_dir)  # pptx-generator/
    image_path = os.path.join(parent_dir, "images", "Arte_metodologia.png")
    
    # Adicionar imagem no centro-direita do slide (maior e mais acima)
    add_local_image(
        slide, 
        image_path, 
        left=4.5,      # Posição à direita
        top=0.2,       # posiçao à abaixo
        width=5.3,     # Largura 
        height=4.8     # Altura
    )
    
    # === TEXTOS DOS 5 PILARES === #
    
    # 1. TECNOLOGIA (Topo - acima do círculo)
    add_texto_pilar(
        slide,
        "TECNOLOGIA",
        "Controle de acesso, Barreiras físicas, Iluminação, CFTV, Comunicação",
        left=7.5,
        top=0.2,
        width=2.0,
        align=PP_ALIGN.LEFT,
        cor_titulo=(0, 51, 102)
    )
    
    # 2. INFORMAÇÕES (Direita - ao lado direito do círculo)
    add_texto_pilar(
        slide,
        "INFORMAÇÕES",
        "Ambiente, Ameaças e Histórico de clientes",
        left=8.8,    # Bem à direita
        top=2.0,     # Centralizado verticalmente
        width=1.0,   # Largura menor pois está no canto
        align=PP_ALIGN.LEFT,
        cor_titulo=(0, 180, 170)
    )
    
    # 3. GESTÃO (Inferior Direita)
    add_texto_pilar(
        slide,
        "GESTÃO",
        "Manutenção, Revisão e Atualização",
        left=7.5,    # Direita inferior
        top=4.5,     # Parte inferior
        width=1.5,
        align=PP_ALIGN.LEFT,
        cor_titulo=(102, 178, 255)
    )
    
    # 4. PROCESSOS (Inferior Esquerda)
    add_texto_pilar(
        slide,
        "PROCESSOS",
        "Documentação, Procedimentos, Supressão vegetal",
        left=5.2,    # Esquerda inferior
        top=4.4,     # Parte inferior
        width=1.8,
        align=PP_ALIGN.LEFT,
        cor_titulo=(0, 102, 204)
    )
    
    # 5. PESSOAS (Esquerda)
    add_texto_pilar(
        slide,
        "PESSOAS",
        "Recursos Humanos, Efetivo de segurança e Treinamentos",
        left=4.3,    # À esquerda do círculo
        top=2.2,     # Centralizado verticalmente
        width=1.5,
        align=PP_ALIGN.LEFT,
        cor_titulo=(0, 76, 153)
    )
    
    print("Slide de Metodologia criado com sucesso!", file=sys.stderr)
    return slide