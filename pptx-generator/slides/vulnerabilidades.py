#!/usr/bin/env python3
# -*- coding: utf-8 -*-
import sys
import os
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import base64
from io import BytesIO

# Importar funções de processamento
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..'))
from utils.processar_dados import reorganizar_por_criticidade, mapear_icone_pilar, mapear_cor_criticidade, nomear_nivel_vulnerabilidade


def add_local_image(slide, image_path, left, top, width, height):
    """Adiciona imagem local ao slide"""
    try:
        if os.path.exists(image_path):
            slide.shapes.add_picture(
                image_path, 
                Inches(left), 
                Inches(top),
                width=Inches(width), 
                height=Inches(height)
            )
            return True
        else:
            print(f"Imagem não encontrada: {image_path}", file=sys.stderr, flush=True)
            return False
    except Exception as e:
        print(f"Erro ao adicionar imagem local: {e}", file=sys.stderr, flush=True)
        return False


def criar_slide_vulnerabilidades(pres, dados, itens_slide):

    slide = pres.slides.add_slide(pres.slide_layouts[6])

    # === FUNDO BRANCO === #
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # === PARALELOGRAMO AZUL === #
    para_azul = slide.shapes.add_shape(
        MSO_SHAPE.PARALLELOGRAM,
        Inches(0.2), Inches(0.2), Inches(0.5), Inches(0.5)
    )
    para_azul.fill.solid()
    para_azul.fill.fore_color.rgb = RGBColor(30, 115, 190)
    para_azul.line.fill.background()

    tf_azul = para_azul.text_frame
    tf_azul.text = "7"
    p = tf_azul.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    tf_azul.vertical_anchor = MSO_ANCHOR.MIDDLE

    # === PARALELOGRAMO CINZA === #
    para_cinza = slide.shapes.add_shape(
        MSO_SHAPE.PARALLELOGRAM,
        Inches(0.85), Inches(0.2), Inches(5.0), Inches(0.5)
    )
    para_cinza.fill.solid()
    para_cinza.fill.fore_color.rgb = RGBColor(70, 70, 70)
    para_cinza.line.fill.background()

    tf_cinza = para_cinza.text_frame
    tf_cinza.text = "VULNERABILIDADES"
    tf_cinza.margin_left = Inches(0.2)
    tf_cinza.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_cinza.paragraphs[0].font.bold = True
    tf_cinza.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # === SUBTÍTULO === #
    subtitulo = slide.shapes.add_textbox(
        Inches(2.5), Inches(0.75), Inches(5.0), Inches(0.35)
    )
    p = subtitulo.text_frame.paragraphs[0]
    p.text = "Não Conformidades"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 38, 77)
    p.alignment = PP_ALIGN.CENTER

    # ================= TABELA ================= #
    rows = 1 + len(itens_slide)
    cols = 7

    table = slide.shapes.add_table(
        rows, cols,
        Inches(0.3), Inches(1.2),
        Inches(9.0),
        Inches(0.2 + 0.6 * len(itens_slide))
    ).table

    # Larguras
    widths = [1.0, 0.5, 2.0, 2.0, 1.2, 1.3, 1.0]
    for i, w in enumerate(widths):
        table.columns[i].width = Inches(w)

    # Alturas
    table.rows[0].height = Inches(0.3)
    for i in range(1, rows):
        table.rows[i].height = Inches(0.6)

    # Cabeçalho
    headers = [
        "Elementos", "NC", "Não Conformidade",
        "Observação", "Risco", "Localização", "Criticidade"
    ]

    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(30, 115, 190)

        p = cell.text_frame.paragraphs[0]
        p.text = header
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # ================= RENDERIZAR LINHAS DA TABELA ================= #
    # Caminho das imagens
    current_dir = os.path.dirname(__file__)
    parent_dir = os.path.dirname(current_dir)
    images_dir = os.path.join(parent_dir, "images")
    
    # Guardar posição da tabela para calcular posições das células
    table_left = Inches(0.3)
    table_top = Inches(1.2)
    row_height = Inches(0.6)
    header_height = Inches(0.3)
    
    for row_idx, item in enumerate(itens_slide, start=1):
        # Calcular posição Y desta linha
        cell_top = table_top + header_height + (row_idx - 1) * row_height

        if row_idx % 2 == 0:
            cor_linha = RGBColor(235, 241, 246)  # azul
        else:
            cor_linha = RGBColor(220, 220, 220) # cinza

        for col in range(cols):
            cell = table.cell(row_idx, col)
            cell.fill.solid()
            cell.fill.fore_color.rgb = cor_linha
        
        # Coluna 0: ÍCONE DO ELEMENTO (pilar)
        cell = table.cell(row_idx, 0)
        
        # Carregar o ícone
        icone_nome = item["elemento"]  # Já vem como "ICON_PESSOAS.png" etc
        image_path = os.path.join(images_dir, icone_nome)
        
        # Calcular posição X da coluna 0
        col_0_left = table_left
        
        try:
            if os.path.exists(image_path):
                slide.shapes.add_picture(
                    image_path,
                    col_0_left + Inches(0.25),
                    cell_top + Inches(0.10),
                    width=Inches(0.4)
                )
        except Exception as e:
            print(f"Erro ao carregar ícone {icone_nome}: {e}", file=sys.stderr)
        
        # Coluna 1: NC
        cell = table.cell(row_idx, 1)
        p = cell.text_frame.paragraphs[0]
        p.text = item["nc"]
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(30, 115, 190)
        p.alignment = PP_ALIGN.CENTER
        cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Coluna 2: NÃO CONFORMIDADE
        cell = table.cell(row_idx, 2)
        p = cell.text_frame.paragraphs[0]
        p.text = item["nao_conformidade"]
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(30, 115, 190)
        p.alignment = PP_ALIGN.CENTER
        cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.text_frame.word_wrap = True
        
        # Coluna 3: OBSERVAÇÃO (vazia)
        cell = table.cell(row_idx, 3)
        p = cell.text_frame.paragraphs[0]
        p.text = item["observacao"]
        p.font.size = Pt(10)
        p.alignment = PP_ALIGN.CENTER
        cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Coluna 4: RISCO (vazia)
        cell = table.cell(row_idx, 4)
        p = cell.text_frame.paragraphs[0]
        p.text = item["risco"]
        p.font.size = Pt(10)
        p.alignment = PP_ALIGN.CENTER
        cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Coluna 5: LOCALIZAÇÃO (vazia)
        cell = table.cell(row_idx, 5)
        p = cell.text_frame.paragraphs[0]
        p.text = item["localizacao"]
        p.font.size = Pt(10)
        p.alignment = PP_ALIGN.CENTER
        cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Coluna 6: CRITICIDADE (círculo colorido)
        cell = table.cell(row_idx, 6)
        
        # Adicionar círculo colorido
        criticidade = item["criticidade"]
        cor_rgb = mapear_cor_criticidade(criticidade)
        
        # Calcular posição X da coluna 6 (última coluna)
        col_6_left = table_left + Inches(8.0)
        
        circulo = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            col_6_left + Inches(0.35),
            cell_top + Inches(0.17),
            Inches(0.25),
            Inches(0.25)
        )
        circulo.fill.solid()
        circulo.fill.fore_color.rgb = RGBColor(*cor_rgb)
        circulo.line.fill.background()

    
    # -------- TÍTULO DA LEGENDA  -------- #
    titulo_legenda = slide.shapes.add_textbox(
        Inches(0.3),
        Inches(4.55),   
        Inches(4.2),
        Inches(0.25)
    )

    p = titulo_legenda.text_frame.paragraphs[0]
    p.text = "Criticidade Risco"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.LEFT

    legenda = slide.shapes.add_textbox(
        Inches(0.2),        # Left
        Inches(4.8 ),        # Top
        Inches(4.2 ),        # Width
        Inches(0.4)         # Height 
    )

    # Borda azul escura (mais fina)
    legenda.line.color.rgb = RGBColor(0, 51, 102)
    legenda.line.width = Pt(0.9)

    tf_legenda = legenda.text_frame
    tf_legenda.clear()
    tf_legenda.margin_left = Inches(0.15)
    tf_legenda.margin_top = Inches(0.05)


    # -------- ITENS (HORIZONTAIS) -------- #
    legenda_itens = [
        ("Muito Alto",  RGBColor(192, 0, 0)),
        ("Alto",        RGBColor(237, 125, 49)),
        ("Médio",       RGBColor(255, 192, 0)),
        ("Baixo",       RGBColor(0, 97, 0)),
        ("Muito Baixo", RGBColor(146, 208, 80)),
    ]

    start_x = Inches(0.3)
    y = Inches(4.9)
    espacamento = Inches(0.8)

    for texto, cor in legenda_itens:
        # Círculo grande
        circulo = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            start_x, y,
            Inches(0.18), Inches(0.18)
        )
        circulo.fill.solid()
        circulo.fill.fore_color.rgb = cor
        circulo.line.fill.background()

        # Texto ao lado do círculo
        label = slide.shapes.add_textbox(
            start_x + Inches(0.12),
            y - Inches(0.02),
            Inches(0.6),
            Inches(0.25)
        )   
        p = label.text_frame.paragraphs[0]
        p.text = texto
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(30, 115, 190)
        p.alignment = PP_ALIGN.LEFT

        start_x += espacamento
      
    titulo_elementos = slide.shapes.add_textbox(
        Inches(4.7),     # à direita da legenda anterior
        Inches(4.55),
        Inches(4.8),
        Inches(0.25)
    )

    p = titulo_elementos.text_frame.paragraphs[0]
    p.text = "Elementos"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.LEFT


    # ========= CAIXA DA LEGENDA ========= #
    legenda_elementos = slide.shapes.add_textbox(
        Inches(4.5),     # Left
        Inches(4.8),     # Top
        Inches(4.6),     # Width
        Inches(0.4)      # Height
    )

    
    legenda_elementos.fill.solid()
    legenda_elementos.fill.fore_color.rgb = RGBColor(255,255,255)

    # Borda
    legenda_elementos.line.color.rgb = RGBColor(0, 51, 102)
    legenda_elementos.line.width = Pt(0.9)


    # ========= ITENS DA LEGENDA ========= #
    elementos = [
        ("tecnologia.png",   "Tecnologia"),
        ("processos.png",    "Processos"),
        ("pessoas.png",      "Pessoas"),
        ("informacoes.png",  "Informação"),
        ("gestao.png",       "Gestão"),
    ]

    start_x = Inches(4.6)
    y = Inches(4.85)
    espacamento = Inches(0.95)

    for nome_icone, texto in elementos:
        image_path = os.path.join(images_dir, nome_icone)

        # Ícone
        try:
            slide.shapes.add_picture(
                image_path,
                start_x,
                y,
                width=Inches(0.28)
            )
        except Exception as e:
            print(f"Erro ao carregar ícone {nome_icone}: {e}", file=sys.stderr)

        # Texto ao lado do ícone
        label = slide.shapes.add_textbox(
            start_x + Inches(0.22),
            y - Inches(0.02),
            Inches(0.8),
            Inches(0.3)
        )

        p = label.text_frame.paragraphs[0]
        p.text = texto
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(30, 115, 190)
        p.alignment = PP_ALIGN.LEFT

        start_x += espacamento
    
    # === LOGO === #
    logo = dados.get("imagens", {}).get("logo_empresa")
    if logo:
        try:
            if ',' in logo:
                logo = logo.split(',')[1]
            stream = BytesIO(base64.b64decode(logo))
            slide.shapes.add_picture(stream, Inches(9.1), Inches(4.5), width=Inches(0.8))
        except Exception as e:
            print(f"Erro ao adicionar logo: {e}", file=sys.stderr)

    return slide


def gerar_vulnerabilidades(pres, dados):
    """
    Gera slides de vulnerabilidades com dados reais
    """
    print("🔄 Processando vulnerabilidades...", file=sys.stderr, flush=True)
    
    # Puxar dados reais das respostas
    respostas = dados.get('dados_modelo', {}).get('respostas', [])
    
    if not respostas or len(respostas) == 0:
        print("⚠️ Nenhuma resposta encontrada", file=sys.stderr, flush=True)
        return
    
    # Reorganizar por criticidade e filtrar vulnerabilidades
    vulnerabilidades = reorganizar_por_criticidade(respostas)
    
    if len(vulnerabilidades) == 0:
        print("⚠️ Nenhuma vulnerabilidade encontrada (todas com nível 1)", file=sys.stderr, flush=True)
        return
    
    print(f"📊 {len(vulnerabilidades)} vulnerabilidades encontradas", file=sys.stderr, flush=True)
    
    # Converter para o formato esperado pela tabela
    itens = []
    for vuln in vulnerabilidades:
        # Extrair nome do pilar para pegar o ícone correto
        pilar_nome = vuln.get('pilar', '')
        icone_nome = mapear_icone_pilar(pilar_nome)
        
        # Montar texto da não conformidade: topicos + " - " + nivel_texto
        nome_topico = vuln.get('topicos', 'Sem descrição')
        #nome_topico = nome_topico.encode('latin1').decode('utf-8')
        vulnerabilidade_nivel = vuln.get('vulnerabilidade', 0)
        nivel_texto = nomear_nivel_vulnerabilidade(vulnerabilidade_nivel)
        nao_conformidade_texto = f"{nome_topico} - {nivel_texto}"
        
        # Montar o item no formato da tabela
        item = {
            "elemento": icone_nome,  # Nome do ícone para carregar depois
            "nc": f"NC-{vuln.get('nc_sequencial', '000')}",
            "nao_conformidade": nao_conformidade_texto,
            "observacao": "",  # Vazio por enquanto
            "risco": "",  # Vazio por enquanto
            "localizacao": "",  # Vazio por enquanto
            "criticidade": vuln.get('criticidade', 'amarelo')  # Para mapear a cor do círculo
        }
        itens.append(item)
    
    # Paginar em grupos de 5
    itens_por_slide = 5
    
    for i in range(0, len(itens), itens_por_slide):
        lote = itens[i:i + itens_por_slide]
        criar_slide_vulnerabilidades(pres, dados, lote)
    
    print(f"✅ {len(range(0, len(itens), itens_por_slide))} slides de Vulnerabilidades criados!", file=sys.stderr, flush=True)