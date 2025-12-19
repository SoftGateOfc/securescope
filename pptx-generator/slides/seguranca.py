#!/usr/bin/env python3
# -*- coding: utf-8 -*-
import sys
import os
import base64
from io import BytesIO
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


def gerar_seguranca(pres, dados):
    """Gera o slide de Segurança Pública"""
    print("Criando slide de Segurança Pública...", file=sys.stderr, flush=True)
    slide = pres.slides.add_slide(pres.slide_layouts[6])

    # === FUNDO BRANCO === #
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # === PARALELOGRAMO AZUL (número 4) === #
   
    para_azul = slide.shapes.add_shape(
    MSO_SHAPE.PARALLELOGRAM,
    Inches(0.2),    # Left
    Inches(0.20),   # Top
    Inches(0.5),    # Width
    Inches(0.5)     # Height
)
    para_azul.fill.solid()
    para_azul.fill.fore_color.rgb = RGBColor(30, 115, 190)  # Azul
    para_azul.line.fill.background()

    tf_azul = para_azul.text_frame
    tf_azul.text = "4" 
    p_azul = tf_azul.paragraphs[0]
    p_azul.font.name = "Arial"
    p_azul.font.size = Pt(20)
    p_azul.font.bold = True
    p_azul.font.color.rgb = RGBColor(255, 255, 255)
    p_azul.alignment = PP_ALIGN.CENTER
    tf_azul.vertical_anchor = MSO_ANCHOR.MIDDLE

    # === PARALELOGRAMO CINZA (SEGURANÇA PÚBLICA) === #
    # === PARALELOGRAMO CINZA (título) === #
    para_cinza = slide.shapes.add_shape(
    MSO_SHAPE.PARALLELOGRAM,
    Inches(0.85),   # Left
    Inches(0.20),   # Top
    Inches(3.5),    # Width
    Inches(0.5)     # Height
)
    para_cinza.fill.solid()
    para_cinza.fill.fore_color.rgb = RGBColor(70, 70, 70)  # Cinza
    para_cinza.line.fill.background()

    tf_cinza = para_cinza.text_frame
    tf_cinza.text = "SEGURANÇA PUBLICA"  
    tf_cinza.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_cinza.margin_left = Inches(0.2)
    p_cinza = tf_cinza.paragraphs[0]
    p_cinza.alignment = PP_ALIGN.LEFT
    p_cinza.font.name = "Arial"
    p_cinza.font.size = Pt(14)
    p_cinza.font.bold = True
    p_cinza.font.color.rgb = RGBColor(255, 255, 255)

    # === IMAGEM SEGURANÇA PÚBLICA === #
    current_dir = os.path.dirname(__file__)
    parent_dir = os.path.dirname(current_dir)
    imagem_path = os.path.join(parent_dir, "images", "seguranca.png")

    try:
        if os.path.exists(imagem_path):
            slide.shapes.add_picture(
                imagem_path,
                Inches(0.3),      # Margem esquerda
                Inches(1.35),     # Abaixo do título
                width=Inches(8.7),
                height=Inches(3.8)
            )
        else:
            print(f"Imagem não encontrada: {imagem_path}", file=sys.stderr, flush=True)
    except Exception as e:
        print(f"Erro ao adicionar imagem de segurança: {e}", file=sys.stderr, flush=True)

    # === LOGO DA EMPRESA (canto inferior direito) === #
    logo = dados.get("imagens", {}).get("logo_empresa")
    if logo:
        try:
            if ',' in logo:
                logo = logo.split(',')[1]
            
            img_data = base64.b64decode(logo)
            stream = BytesIO(img_data)
            
            slide.shapes.add_picture(
                stream,
                Inches(9.0),   # Canto direito
                Inches(4.5),   # Canto inferior
                width=Inches(0.9)
            )
        except Exception as e:
            print(f"Erro ao adicionar logo: {e}", file=sys.stderr, flush=True)

    print("Slide de Segurança Pública criado!", file=sys.stderr, flush=True)
    return slide