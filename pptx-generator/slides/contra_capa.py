import sys
from io import BytesIO
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from utils.processar_dados import fix_encoding


def add_base64_image(slide, base64_str, left, top, width, height):
    try:
        if ',' in base64_str:
            base64_str = base64_str.split(',')[1]

        import base64
        img_data = base64.b64decode(base64_str)
        stream = BytesIO(img_data)
        slide.shapes.add_picture(
            stream,
            Inches(left),
            Inches(top),
            width=Inches(width),
            height=Inches(height)
        )
        return True
    except Exception as e:
        print(f"Erro imagem: {e}", file=sys.stderr)
        return False


def gerar_contra_capa(pres, dados):
    print("Criando contra capa...", file=sys.stderr)
    slide = pres.slides.add_slide(pres.slide_layouts[6])

    # === FUNDO === #
    try:
        with open('images/CAPA.png', 'rb') as f:
            stream = BytesIO(f.read())
            slide.shapes.add_picture(
                stream,
                Inches(0),
                Inches(0),
                width=pres.slide_width,
                height=pres.slide_height
            )
    except Exception as e:
        print(f"Erro fundo: {e}", file=sys.stderr)

    # === LOGO === #
    logo = dados.get("imagens", {}).get("logo_empresa")
    if logo:
        add_base64_image(slide, logo, 0.2, 0.2, 2.1, 1.2)

    # === BLOCO CENTRAL (mais baixo) === #
    empresa = fix_encoding(
        dados.get("dados", {}).get("nome_empresa", "NOME DA EMPRESA")
    ).upper()

    info_box = slide.shapes.add_textbox(
        Inches(4),
        Inches(3.5),   # ↓ abaixado
        Inches(2),
        Inches(2.2)
    )
    tf = info_box.text_frame
    tf.clear()

    def style(p, size, bold=False, after=4, x= PP_ALIGN.LEFT):
        p.font.name = "Calibri"
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = x
        p.space_before = Pt(0)
        p.space_after = Pt(after)
        p.line_spacing = 0.9

    # === NOME DA EMPRESA === #
    p_emp = tf.paragraphs[0]
    p_emp.text = empresa
    style(p_emp, 18, bold=True, after=6, x=PP_ALIGN.CENTER)

    # === LINHA === #
    linha = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(3.5),
        Inches(4),
        Inches(3.0),
        Inches(0.02)
    )
    linha.fill.solid()
    linha.fill.fore_color.rgb = RGBColor(255, 255, 255)
    linha.line.fill.background()

    # Parágrafo de separação mínimo
    p_sep = tf.add_paragraph()
    p_sep.space_after = Pt(2)

    # === CAMPOS EM UMA LINHA === #
    def campo(texto):
        p = tf.add_paragraph()
        p.text = texto
        style(p, 11, bold=True, after=2, x=PP_ALIGN.LEFT)

    campo("ENDEREÇO: Altere")
    campo("TELEFONE: Altere")
    campo("SITE: Altere")
    campo("REDES SOCIAIS: Altere")

    print("Contra capa criada!", file=sys.stderr)
    return slide    